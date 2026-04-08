"""
Fendt PESTEL-EL Sentinel Dashboard
====================================

Modern Streamlit dashboard powered by SQLite backend (q2_solution).

**Phase 4 Enterprise Features:**
- BLUF AI Executive Narrative (or rule-based fallback)
- Conversational AI Strategy Sparring (or prototype fallback)

Features 6 tabs:
1. Executive Summary: AI/rule-based BLUF + metrics + critical disruptions
2. Innovation Radar: Interactive Plotly visualization with dimension filters
3. Live Signal Feed: Searchable dataframe of all signals
4. The Inquisition: Conversational AI strategy dialogue (or prototype mode)
5. Knowledge Graph: Causal interdependencies (Table View + Network Graph)
6. Strategic Reports: C-suite markdown briefs viewer

**Prototype Mode:** Works without GEMINI_API_KEY using intelligent fallbacks.
**Production Mode:** Configure API key for full Claude AI capabilities.
"""

import streamlit as st
import pandas as pd
from datetime import datetime
from pathlib import Path
import sys
import os
from typing import List, Dict
from dotenv import load_dotenv

# Graceful Gemini import — app must render even if package is absent
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

load_dotenv()

# Add q2_solution to path — multiple strategies for local vs Streamlit Cloud
PROJECT_ROOT = Path(__file__).parent.absolute()
_q2_path = str(PROJECT_ROOT / "q2_solution")
if _q2_path not in sys.path:
    sys.path.insert(0, _q2_path)

try:
    from database import SignalDatabase
except ImportError:
    sys.path.insert(0, str(PROJECT_ROOT))
    from q2_solution.database import SignalDatabase

try:
    from semantic_search import SemanticSearch
    SEMANTIC_SEARCH_AVAILABLE = True
except ImportError:
    try:
        from q2_solution.semantic_search import SemanticSearch
        SEMANTIC_SEARCH_AVAILABLE = True
    except ImportError:
        SEMANTIC_SEARCH_AVAILABLE = False

try:
    from innovation_radar import InnovationRadar, PESTEL_COLORS
except ImportError:
    try:
        from q2_solution.innovation_radar import InnovationRadar, PESTEL_COLORS
    except ImportError:
        from q2_solution.innovation_radar import InnovationRadar
        PESTEL_COLORS = {
            'POLITICAL': '#e41a1c', 'ECONOMIC': '#377eb8', 'SOCIAL': '#4daf4a',
            'TECHNOLOGICAL': '#984ea3', 'ENVIRONMENTAL': '#ff7f00', 'LEGAL': '#ffff33',
            'INNOVATION': '#00ccff', 'SOCIAL_MEDIA': '#ff69b4',
        }

# Preferred Gemini models in priority order (newest stable first)
_GEMINI_MODELS = ["gemini-2.0-flash", "gemini-1.5-flash", "gemini-1.5-flash-8b"]


def get_api_key() -> str:
    """Get Gemini API key from Streamlit secrets (cloud) or env (local)."""
    try:
        if "GEMINI_API_KEY" in st.secrets:
            return st.secrets["GEMINI_API_KEY"]
    except Exception:
        pass
    return os.getenv("GEMINI_API_KEY", "")


def _get_gemini_model():
    """
    Return a configured GenerativeModel or None.

    ZERO probe calls — no API call is made here.
    Quota failures recorded in session state are respected across the entire session.
    """
    if st.session_state.get('_gemini_quota_hit'):
        return None
    api_key = get_api_key()
    if not api_key or not GEMINI_AVAILABLE:
        return None
    import google.generativeai as genai
    genai.configure(api_key=api_key)
    return genai.GenerativeModel(_GEMINI_MODELS[0])


def _call_gemini(model, prompt: str, max_tokens: int = 400) -> str | None:
    """
    Call Gemini and return response text, or None on any failure.

    Catches quota / rate-limit errors and stores them in session state so
    all subsequent calls in the same session skip the API (zero extra quota usage).
    """
    try:
        resp = model.generate_content(
            prompt, generation_config={"max_output_tokens": max_tokens}
        )
        st.session_state['_gemini_status'] = f"OK ({model.model_name})"
        return resp.text.strip()
    except Exception as e:
        err = str(e)
        if "429" in err or "quota" in err.lower() or "RESOURCE_EXHAUSTED" in err:
            st.session_state['_gemini_quota_hit'] = True
            st.session_state['_gemini_status'] = "Daily quota exhausted — rule-based fallback active"
        else:
            st.session_state['_gemini_status'] = f"Gemini error: {err[:80]}"
        return None


def _empty_state(icon: str, title: str, message: str, action: str = "") -> None:
    """Render a stylized, enterprise-grade empty state card."""
    action_html = (
        f"<p style='color:#888;margin:10px 0 0 0;font-size:13px;'>{action}</p>"
        if action else ""
    )
    st.markdown(f"""
    <div style='padding:48px 24px;text-align:center;background:rgba(255,255,255,0.03);
                border:1px dashed rgba(255,255,255,0.12);border-radius:12px;margin:16px 0;'>
        <div style='font-size:52px;margin-bottom:14px;'>{icon}</div>
        <h3 style='color:#cccccc;margin:0 0 8px 0;font-weight:600;'>{title}</h3>
        <p style='color:#888888;margin:0;font-size:15px;'>{message}</p>
        {action_html}
    </div>
    """, unsafe_allow_html=True)


def _api_unavailable_state(feature: str) -> None:
    """Render a stylized warning when Gemini API is unavailable (missing key or quota)."""
    status_msg = st.session_state.get('_gemini_status', '')
    is_quota = "quota" in status_msg.lower() or "exhausted" in status_msg.lower()
    icon = "⏳" if is_quota else "⚙️"
    detail = (
        "Daily free-tier quota exhausted. The rule-based fallback is active. "
        "The quota resets at midnight UTC, or upgrade to a paid Gemini plan for uninterrupted AI."
        if is_quota else
        f"<strong>{feature}</strong> requires a <code>GEMINI_API_KEY</code> in Streamlit secrets "
        f"or your <code>.env</code> file. A rule-based fallback is active in the meantime."
    )
    st.markdown(f"""
    <div style='padding:16px 20px;background:rgba(255,153,0,0.08);
                border-left:4px solid #ff9900;border-radius:8px;margin:12px 0;'>
        <h4 style='color:#ff9900;margin:0 0 4px 0;'>{icon} AI Fallback Mode — {feature}</h4>
        <p style='color:#ccc;margin:0;font-size:13px;'>{detail}</p>
        {f"<p style='color:#888;margin:4px 0 0 0;font-size:12px;'>Status: {status_msg}</p>" if status_msg else ""}
    </div>
    """, unsafe_allow_html=True)

# ===========================
# PAGE CONFIGURATION
# ===========================
st.set_page_config(
    page_title="Fendt PESTEL-EL Sentinel",
    page_icon="🛰️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for premium dark theme
st.markdown("""
<style>
    .main {
        background-color: #0e1117;
    }
    .stApp {
        color: #e0e0e0;
    }
    .metric-card {
        padding: 20px;
        border-radius: 10px;
        background: rgba(255, 255, 255, 0.05);
        border: 1px solid rgba(255, 255, 255, 0.1);
        margin-bottom: 15px;
    }
    .critical-alert {
        padding: 15px;
        border-left: 4px solid #ff4b4b;
        background: rgba(255, 75, 75, 0.1);
        border-radius: 5px;
        margin-bottom: 10px;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        padding: 10px 20px;
        background-color: rgba(255, 255, 255, 0.05);
        border-radius: 5px;
    }
</style>
""", unsafe_allow_html=True)

# ===========================
# DATABASE CONNECTION
# ===========================

@st.cache_resource
def get_db():
    """Get cached database connection."""
    return SignalDatabase()

@st.cache_data(ttl=60)  # Cache for 1 minute
def load_all_signals():
    """Load all signals from database."""
    db = get_db()
    return db.get_all_signals()

@st.cache_data(ttl=60)
def get_db_stats():
    """Get database statistics."""
    db = get_db()
    return db.get_database_stats()

# ===========================
# BLUF AI EXECUTIVE NARRATIVE
# ===========================

def generate_bluf_narrative(db_stats: Dict, signals: List[Dict]) -> str:
    """
    Generate Bottom Line Up Front (BLUF) executive narrative.

    Uses Gemini API if available, otherwise generates intelligent rule-based summary.

    Analyzes database statistics and critical signals to produce a 3-sentence
    hard-hitting strategic summary for C-suite consumption.

    Args:
        db_stats: Database statistics from get_database_stats()
        signals: List of all signals from database

    Returns:
        str: 3-sentence BLUF narrative
    """
    # Guard: empty signal list must never crash — return a clear action prompt
    if not signals:
        return (
            "No disruption signals are currently loaded in the database. "
            "Run the pipeline to generate intelligence: `python sentinel.py --run-once`. "
            "Once signals are ingested, this briefing will update automatically."
        )

    api_key = get_api_key()

    # Analyze signal distribution
    df = pd.DataFrame(signals)

    # Count by dimension — safe even if column has mixed/missing values
    dimension_counts = df['primary_dimension'].dropna().value_counts().to_dict()
    if not dimension_counts:
        top_dimension, top_dimension_pct = "Unknown", 0.0
    else:
        top_dimension = max(dimension_counts, key=dimension_counts.get)
        top_dimension_pct = (dimension_counts[top_dimension] / len(df)) * 100

    # Count by severity — column may not exist in older schema
    severity_col = 'disruption_classification'
    if severity_col in df.columns:
        severity_counts = df[severity_col].dropna().value_counts().to_dict()
    else:
        severity_counts = {}
    critical_count = severity_counts.get('CRITICAL', 0)
    high_count = severity_counts.get('HIGH', 0)

    # Calculate average scores
    avg_velocity = df['velocity_score'].mean() if 'velocity_score' in df.columns else 0

    # Get top critical signal
    critical_df = df[df['disruption_classification'] == 'CRITICAL'].sort_values('impact_score', ascending=False)
    top_threat = critical_df.iloc[0]['title'] if len(critical_df) > 0 else None

    # Use rule-based fallback if Gemini is unavailable or quota exhausted
    gemini_model = _get_gemini_model()
    if not gemini_model:
        # Sentence 1: Overall status and volume
        sentence1 = f"Currently tracking {db_stats['total_signals']} disruption signals with {critical_count} CRITICAL and {high_count} HIGH priority threats demanding immediate attention."

        # Sentence 2: Dimension analysis
        if top_dimension_pct > 40:
            sentence2 = f"Signals are heavily skewed toward {top_dimension} dimension ({top_dimension_pct:.0f}% of total), indicating concentrated regulatory/market pressure in this area."
        else:
            sentence2 = f"Disruptions are distributed across PESTEL dimensions with {top_dimension} showing highest concentration ({top_dimension_pct:.0f}%)."

        # Sentence 3: Velocity and top threat
        if avg_velocity > 0.6:
            velocity_desc = "rapidly accelerating"
        elif avg_velocity > 0.4:
            velocity_desc = "building momentum"
        else:
            velocity_desc = "emerging gradually"

        if top_threat:
            sentence3 = f"Velocity indicators show {velocity_desc} disruption patterns, with most critical threat identified as: {top_threat[:80]}..."
        else:
            sentence3 = f"Velocity indicators show {velocity_desc} disruption patterns across the intelligence landscape."

        return f"{sentence1} {sentence2} {sentence3}"

    # AI path — Gemini is healthy, generate sophisticated narrative
    avg_impact = df['impact_score'].mean() if 'impact_score' in df.columns else 0
    avg_novelty = df['novelty_score'].mean() if 'novelty_score' in df.columns else 0

    context = f"""DATABASE STATISTICS:
- Total Signals: {db_stats['total_signals']}
- Date Range: {db_stats['date_range']['earliest']} to {db_stats['date_range']['latest']}

PESTEL DISTRIBUTION:
{chr(10).join(f'- {dim}: {count} signals ({count/db_stats["total_signals"]*100:.1f}%)' for dim, count in dimension_counts.items())}

SEVERITY BREAKDOWN:
{chr(10).join(f'- {severity}: {count}' for severity, count in severity_counts.items())}

AVERAGE SCORES: Impact {avg_impact:.2f} | Velocity {avg_velocity:.2f} | Novelty {avg_novelty:.2f}

TOP CRITICAL THREAT: {top_threat}"""

    prompt = (
        "You are a strategic advisor delivering a BLUF (Bottom Line Up Front) briefing "
        "to AGCO/Fendt's C-suite. Using ONLY the statistical data provided below — "
        "no external knowledge, no fabricated URLs, no invented company names or figures — "
        "write EXACTLY 3 hard-hitting sentences:\n"
        "Sentence 1: The headline finding — what the data shows right now.\n"
        "Sentence 2: The strategic insight — what this PESTEL pattern means specifically "
        "for AGCO/Fendt's competitive position, revenue, and R&D roadmap.\n"
        "Sentence 3: The single most important C-suite action derived strictly from the data.\n"
        "Rules: No preamble. No titles. No bullet points. No URLs. No fabricated facts. "
        "ONLY the 3 sentences.\n\n"
        f"{context}"
    )

    ai_text = _call_gemini(gemini_model, prompt, max_tokens=300)
    if ai_text:
        return ai_text

    # Fallback if AI call fails mid-flight
    return (
        f"Currently tracking {db_stats['total_signals']} disruption signals with "
        f"{critical_count} CRITICAL and {high_count} HIGH priority threats. "
        f"{top_dimension} dimension shows highest concentration ({top_dimension_pct:.0f}%). "
        f"{'Priority threat: ' + top_threat[:80] + '...' if top_threat else 'No critical threats detected.'}"
    )

# ===========================
# THE INQUISITION - GEMINI API
# ===========================

def generate_strategic_questions(critical_signals: List[Dict]) -> List[str]:
    """
    Generate 3-5 hard-hitting strategic questions for Fendt C-suite
    using Gemini API based on critical signals.

    Args:
        critical_signals: List of critical disruption signals

    Returns:
        List of strategic questions
    """
    _fallback_questions = [
        "If EU electrification mandates accelerate to 2027, which Fendt R&D programs face obsolescence?",
        "Which competitor is closest to commercializing autonomous precision farming, and what is Fendt's countermove?",
        "How does the CAP subsidy reallocation toward sustainability affect Fendt's core diesel tractor revenue?",
    ]

    gemini_model = _get_gemini_model()
    if not gemini_model:
        return _fallback_questions

    signal_summaries = "\n".join(
        f"- {sig['title']} [{sig['primary_dimension']}] "
        f"(Disruption: {sig.get('disruption_classification', 'N/A')}, "
        f"Impact: {sig.get('impact_score', 0):.2f})"
        for sig in critical_signals[:10]
    )

    prompt = (
        "You are a strategic advisor to AGCO/Fendt's senior leadership. "
        "Based on these disruption signals, generate 3-5 hard-hitting strategic questions "
        "numbered 1-5. No preamble. ONLY the questions.\n\n"
        f"SIGNALS:\n{signal_summaries}"
    )

    ai_text = _call_gemini(gemini_model, prompt, max_tokens=400)
    if ai_text:
        questions = []
        for line in ai_text.split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('-')):
                q = line.lstrip('0123456789.-) *').strip()
                if q:
                    questions.append(q)
        return questions if questions else _fallback_questions
    return _fallback_questions

# ===========================
# MAIN DASHBOARD
# ===========================

st.markdown("""
<div style='text-align: center; padding: 20px; background: linear-gradient(135deg, rgba(0,100,0,0.1) 0%, rgba(0,50,0,0.2) 100%); border-radius: 10px; margin-bottom: 20px;'>
    <h1 style='color: #00ff88; margin: 0;'>🛰️ Fendt PESTEL-EL Sentinel</h1>
    <h3 style='color: #00ccff; margin: 10px 0;'>Strategic Intelligence War Room</h3>
    <p style='color: #999; margin: 0;'>Autonomous Agricultural Disruption Detection for AGCO/Fendt C-Suite</p>
</div>
""", unsafe_allow_html=True)

# Sidebar
st.sidebar.title("System Status")
stats = get_db_stats()

if stats['total_signals'] > 0:
    st.sidebar.metric("Total Signals", stats['total_signals'])

    if stats['date_range']:
        st.sidebar.caption(f"Earliest: {stats['date_range']['earliest']}")
        st.sidebar.caption(f"Latest: {stats['date_range']['latest']}")
else:
    st.sidebar.warning("No signals in database yet.")
    st.sidebar.info("Run the daily intelligence sweep to populate data.")

st.sidebar.markdown("---")

# Gemini API status — read from session state only (no probe calls, no quota usage)
_quota_hit = st.session_state.get('_gemini_quota_hit', False)
_has_key = bool(get_api_key()) and GEMINI_AVAILABLE
if _quota_hit:
    _g_color, _g_icon, _g_status = "#ff9933", "🟡", "Quota exhausted — fallback active"
elif _has_key:
    _g_status_stored = st.session_state.get('_gemini_status', '')
    if _g_status_stored.startswith("OK"):
        _g_color, _g_icon, _g_status = "#00ff88", "🟢", _g_status_stored
    else:
        _g_color, _g_icon, _g_status = "#00ccff", "🔵", "Ready (on-demand)"
else:
    _g_color, _g_icon, _g_status = "#888888", "⚫", "No API key"
st.sidebar.markdown(
    f"<span style='font-size:12px;color:{_g_color};'>{_g_icon} AI: {_g_status}</span>",
    unsafe_allow_html=True,
)

st.sidebar.markdown("---")
st.sidebar.markdown("**PESTEL-EL Breakdown**")
if stats['signals_per_dimension']:
    for dim, count in stats['signals_per_dimension'].items():
        color = PESTEL_COLORS.get(dim, '#999999')
        st.sidebar.markdown(
            f"<span style='color:{color};font-size:16px;'>●</span> "
            f"<span style='font-size:13px;'><b>{dim}</b>: {count}</span>",
            unsafe_allow_html=True
        )

# Load signals
signals = load_all_signals()

# ===========================
# TAB LAYOUT
# ===========================

tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs([
    "📊 Intelligence Overview",
    "🎯 Innovation Radar",
    "📡 PESTEL Signal Monitor",
    "⚔️ The Inquisition",
    "🕸️ Knowledge Graph",
    "📄 Strategic Reports",
    "📐 Strategic Intelligence Lens",
])

# ===========================
# TAB 1: EXECUTIVE SUMMARY
# ===========================

with tab1:
    st.subheader("Intelligence Overview")
    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0; font-size: 14px;'>
        This tab provides an automated overview of disruption signals detected, classified, and scored by the Sentinel.
        Signals are ingested from live EU agricultural intelligence sources (EUR-Lex, Eurostat, news APIs) and evaluated
        across eight PESTEL-EL dimensions. Metrics below reflect the current contents of the intelligence database —
        they are updated each time the daily sweep runs.
        </p>
    </div>
    """, unsafe_allow_html=True)

    if not signals:
        _empty_state(
            "📡", "Intelligence Database Empty",
            "No disruption signals have been ingested yet.",
            "Run the daily intelligence sweep to populate the radar: <code>python sentinel.py --run-once</code>"
        )
    else:
        # ===========================
        # BLUF AI EXECUTIVE NARRATIVE
        # ===========================

        st.markdown("""
        <div style='padding: 20px; background: linear-gradient(135deg, rgba(0,100,200,0.15) 0%, rgba(0,50,100,0.2) 100%); border-left: 4px solid #00ccff; border-radius: 10px; margin-bottom: 10px;'>
            <h3 style='color: #00ccff; margin: 0 0 6px 0;'>🎯 AI-Generated Signal Synthesis (BLUF)</h3>
            <p style='color: #888; margin: 0; font-size: 13px;'>
                Bottom Line Up Front — a 3-sentence strategic summary of the signals currently in the intelligence database,
                generated automatically from PESTEL distribution, severity counts, and disruption scores.
                <b>Note:</b> "Novelty" measures how unique each signal is <em>relative to previously ingested signals in this database</em>,
                not whether the event is globally new. A high novelty score means the Sentinel has not seen this type of signal before.
            </p>
        </div>
        """, unsafe_allow_html=True)

        # Generate or retrieve BLUF from session state
        if 'bluf_narrative' not in st.session_state or st.button("🔄 Regenerate AI Narrative", help="Click to generate a fresh AI analysis of current intelligence data"):
            with st.spinner("Generating AI executive narrative..."):
                st.session_state.bluf_narrative = generate_bluf_narrative(stats, signals)

        # Display BLUF with help tooltip
        st.markdown(
            f"""
            <div style='font-size: 18px; line-height: 1.8; color: #e0e0e0; padding: 15px; background: rgba(0,0,0,0.3); border-radius: 8px; margin-bottom: 25px;'>
                {st.session_state.bluf_narrative}
            </div>
            """,
            unsafe_allow_html=True
        )

        if not get_api_key() or not GEMINI_AVAILABLE:
            _api_unavailable_state("BLUF AI Executive Narrative")
        else:
            st.caption("ℹ️ **AI-Generated Synthesis:** Dynamically generated by Gemini analyzing PESTEL distribution, severity breakdown, and mathematical disruption scores.")

        st.markdown("---")

        # Convert to DataFrame for easier manipulation
        df = pd.DataFrame(signals)

        # Top-level metrics
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric(
                "Total Signals", len(signals),
                help="Total number of classified disruption signals ingested from monitored EU agricultural intelligence sources (EUR-Lex, Eurostat, news APIs). Each signal is a detected event or trend that has been scored for Impact, Novelty, and Velocity."
            )

        with col2:
            critical_count = len(df[df['disruption_classification'] == 'CRITICAL'])
            st.metric(
                "🔴 Critical", critical_count,
                help="CRITICAL signals are classified as immediate-horizon threats (0–12 months) with a composite disruption score above the CRITICAL threshold. These require urgent review by senior leadership (CEO, CFO, CTO, VP R&D)."
            )

        with col3:
            high_count = len(df[df['disruption_classification'] == 'HIGH'])
            st.metric(
                "🟠 High", high_count,
                help="HIGH signals are near-term threats (12–24 months) with elevated disruption scores. These require strategic planning and resource allocation decisions at the VP/Director level."
            )

        with col4:
            # Average disruption score
            if 'impact_score' in df.columns and df['impact_score'].notna().any():
                avg_impact = df['impact_score'].mean()
                st.metric("Avg Impact", f"{avg_impact:.2f}", help="Average Impact across all ingested signals. This tracks the overall magnitude of industry disruption currently affecting the European agricultural landscape.")
            else:
                st.metric("Avg Impact", "N/A", help="Average Impact across all ingested signals. This tracks the overall magnitude of industry disruption currently affecting the European agricultural landscape.")

        st.markdown("---")

        # Time-Series Velocity Chart
        st.markdown("### 📈 Temporal Velocity Tracking")
        st.markdown("""
        <div style='padding: 10px 14px; background: rgba(255,255,255,0.03); border-radius: 6px; margin-bottom: 10px; border-left: 3px solid rgba(0,204,255,0.4);'>
            <p style='color: #aaa; margin: 0; font-size: 13px;'>
            <b>What this chart shows:</b> The number of new disruption signals ingested per day over time.
            A rising curve indicates accelerating activity in the intelligence feed (more events detected).
            A flat or declining curve indicates a quieter period. This view enables tracking whether disruption
            momentum is building, stable, or decelerating — distinct from any individual signal's Velocity score
            (which measures a single signal's acceleration relative to historical averages).
            </p>
        </div>
        """, unsafe_allow_html=True)

        # Group signals by date and count
        if 'date_ingested' in df.columns:
            df['date_parsed'] = pd.to_datetime(df['date_ingested']).dt.date
            daily_counts = df.groupby('date_parsed').size().reset_index(name='count')
            daily_counts = daily_counts.sort_values('date_parsed')
            
            # Format explicitly as strings to prevent Plotly from rendering continuous sub-second ticks
            daily_counts['date_str'] = daily_counts['date_parsed'].apply(lambda x: x.strftime('%b %d, %Y'))

            if len(daily_counts) > 0:
                # Create Plotly line chart for better aesthetics
                import plotly.graph_objects as go

                fig = go.Figure()
                fig.add_trace(go.Scatter(
                    x=daily_counts['date_str'],
                    y=daily_counts['count'],
                    mode='lines+markers',
                    name='Daily Signals',
                    line=dict(color='#00ff88', width=3),
                    marker=dict(size=8, color='#00ff88', line=dict(color='white', width=1)),
                    hovertemplate='<b>%{x}</b><br>Signals: %{y}<extra></extra>'
                ))

                fig.update_layout(
                    template='plotly_dark',
                    plot_bgcolor='rgba(0,0,0,0)',
                    paper_bgcolor='rgba(0,0,0,0)',
                    xaxis_title='Date',
                    yaxis_title='Signal Count',
                    height=250,
                    margin=dict(l=0, r=0, t=10, b=0),
                    xaxis=dict(
                        gridcolor='rgba(255,255,255,0.1)',
                        type='category',
                        tickangle=-45
                    ),
                    yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
                    showlegend=False
                )

                st.plotly_chart(fig, use_container_width=True)
            else:
                st.info("Insufficient data for temporal chart. Need signals from multiple days.")
        else:
            st.warning("date_ingested field not found in signals.")

        st.markdown("---")

        # ─── Top 3 Critical Disruptions ────────────────────────────────────
        st.markdown("### 🚨 Top 3 Highest-Priority Signals")
        st.markdown("""
        <div style='padding: 10px 14px; background: rgba(255,75,75,0.06); border-radius: 6px; margin-bottom: 10px; border-left: 3px solid rgba(255,75,75,0.5);'>
            <p style='color: #aaa; margin: 0; font-size: 13px;'>
            <b>Selection criteria:</b> All <span style='color:#ff4b4b;'>CRITICAL</span> signals are considered first (composite disruption score ≥ 0.70),
            then <span style='color:#ff9933;'>HIGH</span> signals (score ≥ 0.50) fill remaining slots.
            Within each tier, signals are ranked by <b>Impact Score</b> — the direct magnitude of disruption
            on Fendt's market, regulatory environment, or supply chain.
            These three signals require immediate C-suite attention.
            </p>
        </div>
        """, unsafe_allow_html=True)

        # Build pool: CRITICAL first, then HIGH; sorted by impact within each tier
        _severity_rank = {'CRITICAL': 0, 'HIGH': 1}
        top_pool = df[df['disruption_classification'].isin(['CRITICAL', 'HIGH'])].copy()
        if not top_pool.empty:
            top_pool['_sev_rank'] = top_pool['disruption_classification'].map(_severity_rank).fillna(2)
            top_pool['_impact'] = top_pool['impact_score'].fillna(0)
            top_pool = top_pool.sort_values(['_sev_rank', '_impact'], ascending=[True, False])
            top_signals = top_pool.drop(columns=['_sev_rank', '_impact']).to_dict('records')
        else:
            top_signals = []

        _rank_emoji = {0: '🥇', 1: '🥈', 2: '🥉'}
        _sev_badge = {'CRITICAL': ('🔴', '#ff4b4b'), 'HIGH': ('🟠', '#ff9933')}

        if top_signals:
            for i, signal in enumerate(top_signals[:3]):
                sev = signal.get('disruption_classification', 'HIGH')
                sev_icon, sev_color = _sev_badge.get(sev, ('⚪', '#aaa'))
                rank = _rank_emoji.get(i, f'#{i+1}')
                impact = signal.get('impact_score') or 0
                novelty = signal.get('novelty_score') or 0
                velocity = signal.get('velocity_score') or 0
                comp_score = round(impact * 0.5 + novelty * 0.3 + velocity * 0.2, 3)

                st.markdown(
                    f"<div style='padding:2px 0 4px 0;'>"
                    f"<span style='font-size:13px;color:#888;'>{rank} Rank {i+1} — "
                    f"Selected because: {sev} severity "
                    f"(composite score {comp_score:.3f}) · Impact {impact:.2f} highest in tier</span>"
                    f"</div>",
                    unsafe_allow_html=True,
                )

                with st.expander(
                    f"{sev_icon} {signal['title']} [{signal['primary_dimension']}]",
                    expanded=(i == 0),
                ):
                    col_a, col_b = st.columns([2, 1])
                    with col_a:
                        st.markdown(f"**Source:** [{signal.get('source', 'N/A')}]({signal.get('url', '#')})")
                        st.markdown(f"**Date:** {signal.get('date_ingested', 'N/A')}")
                        content = signal.get('content', '')
                        st.caption(content[:500] + "…" if len(content) > 500 else content)
                    with col_b:
                        st.metric(
                            "Impact", f"{impact:.2f}",
                            help="Impact (0.0–1.0): Magnitude of disruption, calculated from cross-PESTEL reach and high-leverage triggers (legal mandates, technology breakthroughs, subsidy changes)."
                        )
                        st.metric(
                            "Novelty", f"{novelty:.2f}",
                            help="Novelty (0.0–1.0): Uniqueness vs. existing database entries via text-similarity matching. High novelty = new type of signal not previously seen by the Sentinel."
                        )
                        st.metric(
                            "Velocity", f"{velocity:.2f}",
                            help="Velocity (0.0–1.0): Momentum — 30-day vs. 6-month historical signal-volume ratio. High velocity = topic actively accelerating."
                        )
                        st.metric(
                            "Comp. Score", f"{comp_score:.3f}",
                            help="Composite disruption score = Impact×0.5 + Novelty×0.3 + Velocity×0.2. Drives severity classification: CRITICAL ≥ 0.70 · HIGH ≥ 0.50."
                        )

            remaining = len(top_signals) - 3
            if remaining > 0:
                st.info(f"ℹ️ **{remaining} additional CRITICAL/HIGH signals** detected. View all in the PESTEL Signal Monitor tab.")
        else:
            st.success("✅ No CRITICAL or HIGH disruptions detected. Sentinel is actively monitoring for emerging threats.")

# ===========================
# TAB 2: INNOVATION RADAR
# ===========================

with tab2:
    st.subheader("Innovation Radar — Industry Disruption Map")

    if not signals:
        _empty_state(
            "🎯", "Radar Awaiting Data",
            "No signals available to plot on the disruption radar.",
            "Run the daily intelligence sweep to populate the radar: <code>python sentinel.py --run-once</code>"
        )
    else:
        # ── Legend & time-horizon tooltip ────────────────────────────────
        leg_col, info_col = st.columns([2, 1])
        with leg_col:
            st.markdown("""
            **Ring color = Time Horizon (urgency):**  🔴 **12 Month** — act now  ·  🟡 **24 Month** — pilot phase  ·  🟢 **36 Month** — monitor
            """)
        with info_col:
            st.markdown(
                "<span title='How is the time horizon determined?&#10;&#10;"
                "The Sentinel maps each signal&apos;s composite disruption score "
                "to a time horizon ring:&#10;"
                "• CRITICAL (score ≥ 0.70) → 🔴 12 Month ring. "
                "These signals represent immediate threats — regulatory enforcement, "
                "competitor launches, or subsidy cuts happening within 12 months.&#10;"
                "• HIGH (score ≥ 0.50) → 🟡 24 Month ring. "
                "Signals entering the go/no-go window — begin pilot programs now.&#10;"
                "• MODERATE / LOW (score < 0.50) → 🟢 36 Month ring. "
                "Emerging trends to track; not yet requiring capital allocation.&#10;&#10;"
                "Score formula: Impact×0.5 + Novelty×0.3 + Velocity×0.2'>"
                "ℹ️ <u style='cursor:help;color:#00ccff;'>How is Time Horizon assigned?</u></span>",
                unsafe_allow_html=True,
            )

        # Dimension filter — strict 6 PESTEL; deprecated tags remapped automatically
        _STRICT_PESTEL = {'POLITICAL', 'ECONOMIC', 'SOCIAL', 'TECHNOLOGICAL', 'ENVIRONMENTAL', 'LEGAL'}
        _DEPRECATED_REMAP = {'INNOVATION': 'TECHNOLOGICAL', 'SOCIAL_MEDIA': 'SOCIAL'}

        df_radar = pd.DataFrame(signals)
        # Remap deprecated dim values before filtering so old DB entries still appear
        df_radar['primary_dimension'] = df_radar['primary_dimension'].apply(
            lambda d: _DEPRECATED_REMAP.get(d, d)
        )
        available_dimensions = sorted(
            [d for d in df_radar['primary_dimension'].unique() if d in _STRICT_PESTEL]
        )
        # Warn user if deprecated tags were remapped
        deprecated_found = [s.get('primary_dimension') for s in signals
                            if s.get('primary_dimension') in _DEPRECATED_REMAP]
        if deprecated_found:
            _dep_counts = {}
            for t in deprecated_found:
                _dep_counts[t] = _dep_counts.get(t, 0) + 1
            _dep_msg = ', '.join(f"{t} → {_DEPRECATED_REMAP[t]} ({n})" for t, n in _dep_counts.items())
            st.caption(f"ℹ️ Deprecated dimension tags auto-remapped: {_dep_msg}")

        selected_dimensions = st.multiselect(
            "🎛️ Filter by PESTEL Dimension (strict 6)",
            options=available_dimensions,
            default=available_dimensions,
            help="Standard PESTEL: Political · Economic · Social · Technological · Environmental · Legal. "
                 "Legacy INNOVATION and SOCIAL_MEDIA tags are remapped automatically."
        )

        if not selected_dimensions:
            st.warning("Please select at least one dimension to display the radar.")
        else:
            # Prepare signals for radar — apply deprecated tag remap
            radar_signals = []
            filtered_signals = [
                sig for sig in signals
                if _DEPRECATED_REMAP.get(sig['primary_dimension'], sig['primary_dimension'])
                   in selected_dimensions
            ]

            for sig in filtered_signals:
                classification = sig.get('disruption_classification', 'LOW')
                if classification == 'CRITICAL':
                    horizon = '12_MONTH'
                elif classification == 'HIGH':
                    horizon = '24_MONTH'
                else:
                    horizon = '36_MONTH'

                impact = sig.get('impact_score') or 0
                novelty = sig.get('novelty_score') or 0
                velocity = sig.get('velocity_score') or 0
                disruption_score = impact * 0.5 + novelty * 0.3 + velocity * 0.2

                radar_signals.append({
                    'title': sig['title'],
                    'primary_dimension': _DEPRECATED_REMAP.get(
                        sig['primary_dimension'], sig['primary_dimension']
                    ),
                    'time_horizon': horizon,
                    'disruption_score': disruption_score,
                    'classification': classification,
                    'url': sig.get('url', ''),
                })

            # Generate radar
            try:
                radar = InnovationRadar()
                fig = radar.create_radar(radar_signals)
                st.plotly_chart(fig, use_container_width=True)
            except Exception as e:
                st.error(f"Error generating radar: {str(e)}")
                st.info("Ensure plotly is installed: pip install plotly")

            # ── Dot-Shift History ────────────────────────────────────────
            st.markdown("---")
            st.markdown("#### 🕐 Dot-Shift History — Radar Migration Tracking")

            # Auto-save current snapshot to the DB (idempotent: same date = deduplicated later)
            _snapshot_saved = False
            try:
                _db_for_snapshot = get_db()
                _saved_count = _db_for_snapshot.record_radar_snapshot(radar_signals)
                _snapshot_saved = True
            except Exception:
                pass

            # Check how many distinct snapshot dates are recorded
            try:
                _all_history = get_db().get_all_radar_history()
                _snapshot_dates = sorted({r['snapshot_date'] for r in _all_history})
            except Exception:
                _all_history = []
                _snapshot_dates = []

            if len(_snapshot_dates) >= 2:
                # Identify signals that changed time horizon across runs
                _by_title: Dict[str, List] = {}
                for r in _all_history:
                    _by_title.setdefault(r['signal_title'], []).append(r)

                migrations = []
                for title, history in _by_title.items():
                    horizons = [h['time_horizon'] for h in history]
                    if len(set(horizons)) > 1:
                        migrations.append({
                            'Signal': title[:80],
                            'First Seen': history[0]['snapshot_date'],
                            'First Horizon': history[0]['time_horizon'].replace('_', ' '),
                            'Latest Horizon': history[-1]['time_horizon'].replace('_', ' '),
                            'Moved': '⬆️ Escalated' if horizons[-1] < horizons[0]
                                     else '⬇️ De-escalated',
                        })

                if migrations:
                    st.markdown(
                        f"**{len(migrations)} signal(s) have shifted time-horizon rings** "
                        f"across {len(_snapshot_dates)} tracked pipeline runs:"
                    )
                    st.dataframe(pd.DataFrame(migrations), use_container_width=True, hide_index=True)
                else:
                    st.info(
                        f"No horizon shifts detected yet across {len(_snapshot_dates)} runs. "
                        "Signals have maintained their time-horizon classification."
                    )
            else:
                st.markdown(
                    "<div style='padding:10px 14px;background:rgba(0,204,255,0.07);"
                    "border-radius:6px;border-left:3px solid #00ccff;'>"
                    "<b style='color:#00ccff;'>📊 Dot-Shift Tracking: Active</b><br>"
                    "<span style='color:#aaa;font-size:13px;'>"
                    "This run has been recorded as the baseline snapshot. "
                    "After each subsequent pipeline run, the system will compare positions "
                    "and display arrows here showing signals that have moved from "
                    "36-month → 24-month → 12-month rings (escalation) or the reverse."
                    "</span></div>",
                    unsafe_allow_html=True,
                )

            # Signal table with source links
            st.markdown("---")
            st.markdown("### 📋 Plotted Signals — Source Reference")
            st.caption("All signals currently displayed on the radar, with direct links to source documents.")

            # Build table from original signals (with URL) for filtered dimensions
            radar_table_data = []
            for sig in filtered_signals:
                classification = sig.get('disruption_classification', 'LOW')
                if classification == 'CRITICAL':
                    horizon_label = '🔴 12 Month'
                elif classification == 'HIGH':
                    horizon_label = '🟡 24 Month'
                else:
                    horizon_label = '🟢 36 Month'

                impact = sig.get('impact_score', 0) if sig.get('impact_score') is not None else 0
                novelty = sig.get('novelty_score', 0) if sig.get('novelty_score') is not None else 0
                velocity = sig.get('velocity_score', 0) if sig.get('velocity_score') is not None else 0
                disruption_score = round(impact * 0.5 + novelty * 0.3 + velocity * 0.2, 3)

                radar_table_data.append({
                    'Title': sig['title'],
                    'Dimension': sig['primary_dimension'],
                    'Severity': classification,
                    'Time Horizon': horizon_label,
                    'Disruption Score': disruption_score,
                    'Source URL': sig.get('url', ''),
                })

            radar_table_df = pd.DataFrame(radar_table_data).sort_values('Disruption Score', ascending=False)
            st.dataframe(
                radar_table_df,
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Title": st.column_config.TextColumn("Signal Title", width="large"),
                    "Dimension": st.column_config.TextColumn("PESTEL Dimension", width="small"),
                    "Severity": st.column_config.TextColumn("Severity", width="small"),
                    "Time Horizon": st.column_config.TextColumn("Time Horizon", width="small",
                        help="Horizon is derived from severity: CRITICAL → 12 months, HIGH → 24 months, MODERATE/LOW → 36 months."),
                    "Disruption Score": st.column_config.NumberColumn("Score", format="%.3f", width="small",
                        help="Composite score = Impact×0.5 + Novelty×0.3 + Velocity×0.2"),
                    "Source URL": st.column_config.LinkColumn("🔗 Source", width="medium",
                        help="Direct link to the original source document."),
                }
            )

            # ── Outcome & Innovation Definition ───────────────────────────
            st.markdown("---")
            with st.expander("📖 Radar Outcomes & Innovation Definition — What does this chart tell us?", expanded=False):
                st.markdown("""
                #### What this Radar Measures
                The Innovation Radar maps **strategic disruption signals** across the six PESTEL dimensions
                to help AGCO/Fendt leadership understand *where* competitive forces are building and *when* they
                require a capital allocation decision.

                #### How "Innovation" is Defined Here
                In this framework, **innovation is not limited to patents or product launches**.
                A signal is classified as innovation-relevant when it:
                - Represents a **technology readiness level (TRL) shift** — e.g. a pilot program reaching
                  commercial scale within the tracking period.
                - Involves **R&D investment surges** detected in funding data, patent filings, or
                  academic publication rates.
                - Represents a **business model disruption** — e.g. a new entrant displacing an incumbent's
                  go-to-market with a platform or service model.

                Technological signals are captured under the **TECHNOLOGICAL** PESTEL dimension.
                Regulatory drivers of innovation (e.g. EU mandates forcing electrification) are
                captured under **LEGAL** or **POLITICAL** — because the force is regulatory, not technological.

                #### Derived Findings from the Current Radar
                - **Ring density** (many dots in the 12-month ring) = near-term strategic pressure.
                  Fendt should be in execution mode for these, not discovery mode.
                - **Quadrant concentration** (many signals in one PESTEL sector) = systemic risk in that
                  area — e.g. a cluster of ENVIRONMENTAL signals means regulatory exposure is unusually high.
                - **Score size** (dot radius) = composite disruption magnitude.
                  Large dots require proportionally larger strategic responses.

                #### What the Radar Does NOT Show
                - It does not forecast specific dates — horizons are probabilistic windows.
                - It does not show competitive intelligence about specific rival firms unless
                  those signals are present in the database.
                - Source citations are in the signal table below the radar.
                """)

# ===========================
# TAB 3: LIVE SIGNAL FEED
# ===========================

with tab3:
    st.subheader("PESTEL Signal Monitor")
    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0 0 8px 0; font-size: 14px;'>
        Raw intelligence signals ingested from monitored sources, fully searchable and filterable.
        Each signal has been classified into a PESTEL-EL dimension and scored on three axes:
        </p>
        <ul style='color: #888; margin: 0; font-size: 13px; padding-left: 18px;'>
            <li><b>Impact (0–1):</b> Magnitude of the disruption — how many PESTEL dimensions are affected and whether it involves a high-leverage trigger (legal mandate, technology breakthrough, subsidy change).</li>
            <li><b>Novelty (0–1):</b> Uniqueness relative to signals already in this database — calculated via text similarity matching against historical entries. High novelty = the Sentinel has not seen this type of event before.</li>
            <li><b>Velocity (0–1):</b> Momentum — mathematical comparison of 30-day signal volume vs. the trailing 6-month average. High velocity = this topic is actively accelerating.</li>
        </ul>
        <p style='color: #888; margin: 8px 0 0 0; font-size: 13px;'>
        <b>Severity classification</b> is derived from the composite disruption score
        (Impact×0.5 + Novelty×0.3 + Velocity×0.2): CRITICAL ≥ 0.7 · HIGH ≥ 0.5 · MODERATE ≥ 0.3 · LOW &lt; 0.3.
        Duplicate signals (same source URL) are filtered at display time.
        </p>
    </div>
    """, unsafe_allow_html=True)

    if not signals:
        _empty_state(
            "📡", "Signal Feed Empty",
            "No disruption signals have been ingested yet.",
            "Run the daily intelligence sweep: <code>python sentinel.py --run-once</code>"
        )
    else:
        df = pd.DataFrame(signals)

        # ── Robust deduplication ─────────────────────────────────────────
        # Priority: highest severity first, then latest ingestion date.
        # Step 1: deduplicate by URL (exact match)
        # Step 2: deduplicate by normalised title (case/space-insensitive)
        # This removes both exact re-ingestions and pipeline double-runs.
        _SEV_ORDER = {'CRITICAL': 0, 'HIGH': 1, 'MODERATE': 2, 'LOW': 3}
        df_deduped = df.copy()
        df_deduped['_sev_rank'] = df_deduped['disruption_classification'].map(_SEV_ORDER).fillna(4)
        df_deduped = df_deduped.sort_values(
            ['_sev_rank', 'date_ingested'], ascending=[True, False]
        )
        if 'url' in df_deduped.columns:
            df_deduped = df_deduped.drop_duplicates(subset=['url'], keep='first')
        if 'title' in df_deduped.columns:
            df_deduped['_title_norm'] = df_deduped['title'].str.strip().str.lower()
            df_deduped = df_deduped.drop_duplicates(subset=['_title_norm'], keep='first')
            df_deduped = df_deduped.drop(columns=['_title_norm'])
        df_deduped = df_deduped.drop(columns=['_sev_rank'])
        df_deduped = df_deduped.reset_index(drop=True)

        # Search and filter controls
        col_search, col_filter = st.columns([2, 1])

        with col_search:
            search_term = st.text_input("🔍 Search signals (title, content, source)", "")

        with col_filter:
            dimension_filter = st.multiselect(
                "Filter by dimension",
                options=df_deduped['primary_dimension'].unique().tolist(),
                default=[]
            )

        # Apply filters
        filtered_df = df_deduped.copy()

        if search_term:
            filtered_df = filtered_df[
                filtered_df['title'].str.contains(search_term, case=False, na=False) |
                filtered_df['content'].str.contains(search_term, case=False, na=False) |
                filtered_df['source'].str.contains(search_term, case=False, na=False)
            ]

        if dimension_filter:
            filtered_df = filtered_df[filtered_df['primary_dimension'].isin(dimension_filter)]

        duplicates_removed = len(df) - len(df_deduped)
        caption_parts = [f"Showing {len(filtered_df)} of {len(df_deduped)} unique signals"]
        if duplicates_removed > 0:
            caption_parts.append(
                f"({duplicates_removed} duplicate{'s' if duplicates_removed > 1 else ''} removed "
                "by URL + title deduplication — highest severity kept)"
            )
        st.caption(" ".join(caption_parts))

        # Display configuration — include URL as clickable link
        display_columns = [
            'title', 'primary_dimension', 'disruption_classification',
            'impact_score', 'novelty_score', 'velocity_score',
            'date_ingested', 'source', 'url'
        ]

        # Only show columns that exist
        display_columns = [col for col in display_columns if col in filtered_df.columns]

        # Display dataframe
        st.dataframe(
            filtered_df[display_columns].sort_values('date_ingested', ascending=False),
            use_container_width=True,
            hide_index=True,
            column_config={
                "title": st.column_config.TextColumn(
                    "Title", width="large",
                    help="Signal title as extracted from the source document."
                ),
                "primary_dimension": st.column_config.TextColumn(
                    "Dimension", width="small",
                    help="The primary PESTEL-EL dimension this signal belongs to, determined by the Classifier agent based on signal content and keyword triggers."
                ),
                "disruption_classification": st.column_config.TextColumn(
                    "Severity", width="small",
                    help="Severity classification derived from composite disruption score: CRITICAL (≥0.7), HIGH (≥0.5), MODERATE (≥0.3), LOW (<0.3). Reflects urgency of response, not absolute danger."
                ),
                "impact_score": st.column_config.NumberColumn(
                    "Impact",
                    format="%.2f",
                    help="Impact (0.0–1.0): Magnitude of disruption. Calculated based on cross-PESTEL reach (how many dimensions the signal touches) and high-leverage triggers like legal mandates or technology breakthroughs."
                ),
                "novelty_score": st.column_config.NumberColumn(
                    "Novelty",
                    format="%.2f",
                    help="Novelty (0.0–1.0): Uniqueness vs. existing database entries. Uses text similarity matching — high novelty means the Sentinel has not previously seen this type of signal. Does NOT indicate global newness."
                ),
                "velocity_score": st.column_config.NumberColumn(
                    "Velocity",
                    format="%.2f",
                    help="Velocity (0.0–1.0): Momentum score. Compares 30-day signal count for this topic vs. 6-month historical average. High velocity = trend actively accelerating."
                ),
                "date_ingested": st.column_config.TextColumn(
                    "Date Ingested", width="small",
                    help="Date the signal was detected and ingested by the Sentinel pipeline."
                ),
                "source": st.column_config.TextColumn(
                    "Source", width="small",
                    help="Publication or data provider name (e.g., EUR-Lex, Eurostat, AgriPulse)."
                ),
                "url": st.column_config.LinkColumn(
                    "🔗 Source Link", width="medium",
                    help="Direct link to the original source document. Click to open in browser."
                ),
            }
        )

        # Export option
        st.download_button(
            label="📥 Download as CSV",
            data=filtered_df.to_csv(index=False),
            file_name=f"fendt_signals_{datetime.now().strftime('%Y%m%d')}.csv",
            mime="text/csv"
        )

# ===========================
# TAB 4: THE INQUISITION (CONVERSATIONAL AI)
# ===========================

with tab4:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(255, 75, 75, 0.1); border-left: 4px solid #ff4b4b; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #ff4b4b; margin: 0;'>⚔️ The Inquisition</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Conversational AI Strategy Sparring for C-Suite</p>
    </div>
    """, unsafe_allow_html=True)
    st.caption("**Engage in strategic dialogue with Claude. Ask follow-up questions, challenge assumptions, and explore alternative scenarios based on current disruption intelligence.**")

    if not signals:
        _empty_state(
            "⚔️", "The Inquisition Awaits Intelligence",
            "No signals available. The Inquisition requires live disruption data to interrogate.",
            "Run the daily intelligence sweep first: <code>python sentinel.py --run-once</code>"
        )
    else:
        df = pd.DataFrame(signals)
        # Accept both CRITICAL and HIGH signals for analysis
        high_priority_signals = df[df['disruption_classification'].isin(['CRITICAL', 'HIGH'])].to_dict('records')

        if not high_priority_signals:
            st.warning("No high-priority signals detected. The Inquisition requires at least one CRITICAL or HIGH signal.")
            st.info("Run the daily intelligence sweep to populate data.")
        else:
            critical_count = len(df[df['disruption_classification'] == 'CRITICAL'])
            high_count = len(df[df['disruption_classification'] == 'HIGH'])
            st.markdown(f"**Analyzing {critical_count} CRITICAL and {high_count} HIGH signals ({len(high_priority_signals)} total)...**")

            # Initialize conversation state
            if 'inquisition_messages' not in st.session_state:
                st.session_state.inquisition_messages = []

            if 'inquisition_initialized' not in st.session_state:
                st.session_state.inquisition_initialized = False

            # Show awaiting state on first load — only fire on explicit button click
            if not st.session_state.inquisition_initialized:
                _empty_state(
                    "🔮", "Ready to Interrogate",
                    "Click the button above to generate hard-hitting strategic questions from your live intelligence data.",
                    "The Inquisition analyzes your CRITICAL and HIGH signals and forces C-suite reflection."
                )

            if st.button("🔮 Start New Strategic Session", type="primary"):
                with st.spinner("Consulting Claude for strategic insights..."):
                    questions = generate_strategic_questions(high_priority_signals)

                # Reset conversation
                st.session_state.inquisition_messages = []

                # Add system context
                context_summary = f"I've analyzed {len(high_priority_signals)} high-priority disruption signals ({critical_count} CRITICAL, {high_count} HIGH) across the PESTEL framework for Fendt/AGCO."

                # Format initial message with questions
                initial_message = f"{context_summary}\n\nHere are the critical strategic questions I need you to confront:\n\n"
                for i, question in enumerate(questions, 1):
                    initial_message += f"**{i}. {question}**\n\n"

                initial_message += "How would you like to respond? You can:\n- Answer a specific question\n- Challenge my assumptions\n- Ask for clarification\n- Explore alternative scenarios"

                # Add to conversation history
                st.session_state.inquisition_messages.append({
                    "role": "assistant",
                    "content": initial_message
                })

                st.session_state.inquisition_initialized = True
                st.rerun()

            # Display conversation history
            st.markdown("### 💬 Strategic Dialogue")

            for message in st.session_state.inquisition_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            # Chat input
            if prompt := st.chat_input("Type your response or question (e.g., 'Assume we can't change diesel roadmap until 2028. How does that affect Question 2?')"):
                # Add user message to history
                st.session_state.inquisition_messages.append({
                    "role": "user",
                    "content": prompt
                })

                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)

                # Generate AI response
                with st.chat_message("assistant"):
                    with st.spinner("Generating strategic analysis..."):
                        inq_model = _get_gemini_model()
                        if not inq_model:
                            status = st.session_state.get('_gemini_status', 'No API key configured')
                            st.warning(
                                f"AI unavailable ({status}). The Inquisition requires Gemini to respond. "
                                "Please try again after midnight UTC when the quota resets, "
                                "or configure a paid Gemini API key in Streamlit secrets."
                            )
                        else:
                            try:
                                signal_context = "\n".join(
                                    f"- {sig['title']} [{sig['primary_dimension']}] "
                                    f"(Impact: {sig.get('impact_score', 0):.2f})"
                                    for sig in high_priority_signals[:10]
                                )
                                chat_context = (
                                    "You are a strategic advisor to AGCO/Fendt's senior leadership "
                                    "in a strategic dialogue about agricultural industry disruptions.\n\n"
                                    f"INTELLIGENCE CONTEXT:\n{signal_context}\n\n"
                                )
                                for msg in st.session_state.inquisition_messages[:-1]:
                                    chat_context += f"{msg['role'].upper()}: {msg['content']}\n\n"
                                chat_context += f"USER: {prompt}\n\nRespond strategically:"

                                full_response = _call_gemini(inq_model, chat_context, max_tokens=600)
                                if full_response:
                                    st.markdown(full_response)
                                    st.session_state.inquisition_messages.append({
                                        "role": "assistant",
                                        "content": full_response,
                                    })
                                else:
                                    st.warning("AI response unavailable — quota may be exhausted.")
                            except Exception as e:
                                st.error(f"Error: {str(e)[:200]}")

            # Show signal context
            with st.expander("📋 View High-Priority Signals Being Analyzed"):
                for sig in high_priority_signals[:10]:
                    st.markdown(f"**{sig['title']}** [{sig['primary_dimension']}]")
                    st.caption(f"Impact: {sig.get('impact_score', 0):.2f} | Source: {sig['source']}")
                    st.markdown("---")

# ===========================
# TAB 5: KNOWLEDGE GRAPH
# ===========================

with tab5:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0, 255, 136, 0.1); border-left: 4px solid #00ff88; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00ff88; margin: 0;'>🕸️ Knowledge Graph</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Causal Interdependencies Across PESTEL Dimensions</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0 0 8px 0; font-size: 14px;'><b>How it works — Calculations Explained:</b></p>
        <ul style='color: #888; margin: 0; font-size: 13px; padding-left: 18px;'>
            <li><b>Nodes</b> represent named entities extracted from intelligence signals — a regulatory policy, a market shift, a technology, or a physical event. Node color = PESTEL dimension (red=Political, blue=Economic, green=Social, purple=Technological, orange=Environmental, amber=Legal).</li>
            <li><b>Edges</b> represent causal relationships between two entities, identified by the Analyst agent from signal content and verbatim quotes. Edge <b>weight</b> range −1.0 to +1.0: positive = one event accelerates or enables the other; negative = one event restricts or counters it.</li>
            <li><b>Distance between nodes</b> is determined by the Barnes-Hut physics simulation: nodes with strong causal links (|weight| close to 1.0) are pulled closer together by stronger spring forces. Weakly linked nodes drift further apart.</li>
            <li><b>Temporal decay:</b> Edge weights decay with a 90-day half-life — older, unconfirmed relationships automatically lose influence. Edges below 0.05 are pruned from the graph entirely.</li>
            <li><b>Provenance:</b> Every edge must have a verified <code>source_url</code> and an <code>exact_quote</code> (min 10 characters) — required for EU Data Act 2026 compliance.</li>
            <li><b>Edge colors:</b> 🟢 Strong positive (>0.5) · 🔵 Moderate positive · 🟠 Moderate negative · 🔴 Strong negative (&lt;−0.5)</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)

    graph_path = Path('./data/graph.json')

    if not graph_path.exists():
        _empty_state(
            "🕸️", "Knowledge Graph Not Yet Initialized",
            "No causal relationships have been mapped between PESTEL entities.",
            "The Graph is populated by the Analyst agents during pipeline execution: <code>python sentinel.py --run-once</code>"
        )
    else:
        try:
            import json
            import networkx as nx

            # Load graph
            with open(graph_path, 'r') as f:
                graph_data = json.load(f)

            # Mirror links/edges to prevent NetworkX version KeyErrors across environments
            if 'links' in graph_data and 'edges' not in graph_data:
                graph_data['edges'] = graph_data['links']
            elif 'edges' in graph_data and 'links' not in graph_data:
                graph_data['links'] = graph_data['edges']

            G = nx.node_link_graph(graph_data)

            # Display stats
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Nodes (Entities)", G.number_of_nodes(), help="Nodes represent distinct entities derived from the intelligence feed. A node can be a regulatory policy, an economic shift, a specific technology, or a physical environment change. The more nodes tracked, the broader our horizon scanning.")
            with col2:
                st.metric("Edges (Relationships)", G.number_of_edges(), help="Edges represent the mathematical links connecting two distinct nodes. This explicitly tracks how one event (e.g., a Political change) ripples out to trigger another event (e.g., an Economic cost).")
            with col3:
                if G.number_of_edges() > 0:
                    avg_weight = sum([attrs.get('weight', 0) for _, _, attrs in G.edges(data=True)]) / G.number_of_edges()
                    st.metric("Avg Edge Weight", f"{avg_weight:.2f}", help="Weight (-1.0 to 1.0) measures the strength of the causal ripple. Positive values accelerate or amplify a trend. Negative values decelerate, regulate, or restrict a trend.")
                else:
                    st.metric("Avg Edge Weight", "N/A", help="Weight (-1.0 to 1.0) measures the strength of the causal ripple. Positive values accelerate or amplify a trend. Negative values decelerate, regulate, or restrict a trend.")

            st.markdown("---")

            # View selection
            view_mode = st.radio(
                "📊 Visualization Mode",
                ["📋 Table View (Recommended)", "🕸️ Network Graph"],
                help="Table View is clearer for large graphs. Network Graph shows spatial relationships."
            )

            if view_mode == "📋 Table View (Recommended)":
                st.markdown("### Causal Relationships")
                st.caption("**Each row shows one causal link. Read as: Source → Relationship → Target**")

                causality_data = []
                for source, target, attrs in G.edges(data=True):
                    source_label = G.nodes[source].get('label', source)
                    source_cat = G.nodes[source].get('category', 'N/A')
                    target_label = G.nodes[target].get('label', target)
                    target_cat = G.nodes[target].get('category', 'N/A')

                    relationship = attrs.get('relationship', 'RELATES_TO')
                    weight = attrs.get('weight', 0.0)
                    quote = attrs.get('exact_quote', 'N/A')
                    source_url = attrs.get('source_url', 'N/A')

                    causality_data.append({
                        'Source': f"{source_label} [{source_cat}]",
                        'Relationship': relationship,
                        'Target': f"{target_label} [{target_cat}]",
                        'Weight': weight,
                        'Effect': '🟢 Positive' if weight > 0 else '🔴 Negative' if weight < 0 else '⚪ Neutral',
                        'Evidence Quote': quote[:100] + '...' if len(quote) > 100 else quote,
                        'Source URL': source_url
                    })

                causality_df = pd.DataFrame(causality_data)
                causality_df['abs_weight'] = causality_df['Weight'].abs()
                causality_df = causality_df.sort_values('abs_weight', ascending=False).drop('abs_weight', axis=1)

                st.dataframe(
                    causality_df,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Source": st.column_config.TextColumn("📍 Source Node", width="medium", help="The inciting PESTEL entity that triggered the disruption."),
                        "Relationship": st.column_config.TextColumn("⚡ Causal Verb", width="small", help="The physical action explaining how the source affects the target."),
                        "Target": st.column_config.TextColumn("🎯 Target Node", width="medium", help="The resulting PESTEL entity that is being impacted."),
                        "Weight": st.column_config.NumberColumn("📊 Weight", format="%.2f", width="small", help="The mathematical strength (-1.0 to 1.0) of the ripple effect."),
                        "Effect": st.column_config.TextColumn("↗️ Effect Type", width="small", help="Whether the trend is being accelerated (positive) or restricted (negative)."),
                        "Evidence Quote": st.column_config.TextColumn("📜 Quote", width="large", help="Exact text pulled from the intelligence agent justifying this link."),
                        "Source URL": st.column_config.LinkColumn("🔗 Source", width="medium", help="Direct link to the source document for verification.")
                    }
                )

                st.download_button(
                    label="📥 Download Causality Table as CSV",
                    data=causality_df.to_csv(index=False),
                    file_name=f"knowledge_graph_causality_{datetime.now().strftime('%Y%m%d')}.csv",
                    mime="text/csv"
                )

            else:
                from pyvis.network import Network
                import tempfile
                import streamlit.components.v1 as components

                st.markdown("### Interactive Network Visualization")
                st.caption("**Drag nodes to explore. Click for details. Arrows show causality direction.**")

                col_filter1, col_filter2 = st.columns(2)

                with col_filter1:
                    min_weight = st.slider(
                        "Minimum Edge Weight (Strength)",
                        min_value=0.0,
                        max_value=1.0,
                        value=0.1,
                        step=0.1,
                        help="Filter out weak relationships to reduce clutter"
                    )

                with col_filter2:
                    all_categories = set()
                    for _, attrs in G.nodes(data=True):
                        all_categories.add(attrs.get('category', 'UNKNOWN'))

                    selected_categories = st.multiselect(
                        "Filter by PESTEL Dimension",
                        options=sorted(all_categories),
                        default=sorted(all_categories),
                        help="Show only nodes from selected dimensions"
                    )

                net = Network(
                    height='650px',
                    width='100%',
                    bgcolor='#0e1117',
                    font_color='#ffffff',
                    directed=True
                )

                net.barnes_hut(
                    gravity=-10000,
                    central_gravity=0.1,
                    spring_length=250,
                    spring_strength=0.001,
                    damping=0.5
                )

                category_colors = PESTEL_COLORS  # Single source of truth from innovation_radar.py

                visible_nodes = set()
                for node_id, attrs in G.nodes(data=True):
                    category = attrs.get('category', 'UNKNOWN')

                    if category not in selected_categories:
                        continue

                    label = attrs.get('label', node_id)
                    color = category_colors.get(category, '#999999')

                    display_label = label[:30] + '...' if len(label) > 30 else label

                    net.add_node(
                        node_id,
                        label=display_label,
                        title=f"<b>{label}</b><br>{category}",
                        color=color,
                        size=35,
                        font={'size': 18, 'color': '#ffffff', 'face': 'arial', 'bold': True}
                    )
                    visible_nodes.add(node_id)

                for source, target, attrs in G.edges(data=True):
                    if source not in visible_nodes or target not in visible_nodes:
                        continue

                    weight = attrs.get('weight', 0.0)

                    if abs(weight) < min_weight:
                        continue

                    relationship = attrs.get('relationship', 'RELATES_TO')

                    if weight > 0.5:
                        edge_color = '#00ff88'
                    elif weight > 0:
                        edge_color = '#00ccff'
                    elif weight < -0.5:
                        edge_color = '#ff0066'
                    elif weight < 0:
                        edge_color = '#ff9933'
                    else:
                        edge_color = '#999999'

                    net.add_edge(
                        source,
                        target,
                        title=f"{relationship} (weight: {weight:.2f})",
                        color=edge_color,
                        width=abs(weight) * 5,
                        arrows={'to': {'enabled': True, 'scaleFactor': 1.5}}
                    )

                with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.html') as f:
                    net.save_graph(f.name)
                    with open(f.name, 'r') as html_file:
                        html_content = html_file.read()

                    components.html(html_content, height=670, scrolling=False)

                st.markdown("---")
                st.markdown("### 🎨 Edge Color Legend")
                col_l1, col_l2, col_l3, col_l4 = st.columns(4)

                with col_l1:
                    st.markdown("🟢 **Strong Positive** (Weight > 0.5)")
                with col_l2:
                    st.markdown("🔵 **Moderate Positive** (Weight 0 to 0.5)")
                with col_l3:
                    st.markdown("🟠 **Moderate Negative** (Weight -0.5 to 0)")
                with col_l4:
                    st.markdown("🔴 **Strong Negative** (Weight < -0.5)")

            # ── AI Forecast Extension (inside the main KG try: block) ────
            st.markdown("---")
            st.markdown("### 🔮 AI Causal Forecast — Cross-PESTEL Predictions")
            st.markdown("""
            <div style='padding: 10px 14px; background: rgba(150,78,163,0.08); border-radius: 6px;
                        border-left: 3px solid #984ea3; margin-bottom: 10px;'>
                <p style='color: #aaa; margin: 0; font-size: 13px;'>
                Based on the verified causal edges currently in the Knowledge Graph, this section uses AI to
                predict <b>next-level interdependencies</b> — how an existing signal in one PESTEL area
                will logically propagate to another area. Forecasts are grounded strictly in the
                existing graph edges; no external sources or fabricated links are introduced.
                </p>
            </div>
            """, unsafe_allow_html=True)

            if st.button("🔮 Generate Cross-PESTEL Forecast", key="kg_forecast_btn"):
                _forecast_model = _get_gemini_model()
                if not _forecast_model:
                    _api_unavailable_state("Cross-PESTEL Forecast")
                else:
                    try:
                        # Build edge context — strictly from existing graph, no invention
                        _edge_lines = []
                        for _src, _tgt, _eattrs in list(G.edges(data=True))[:20]:
                            _src_lbl = G.nodes[_src].get('label', _src)
                            _src_cat = G.nodes[_src].get('category', 'UNKNOWN')
                            _tgt_lbl = G.nodes[_tgt].get('label', _tgt)
                            _tgt_cat = G.nodes[_tgt].get('category', 'UNKNOWN')
                            _wt = _eattrs.get('weight', 0)
                            _rel = _eattrs.get('relationship', 'RELATES_TO')
                            _edge_lines.append(
                                f"  [{_src_cat}] {_src_lbl} --{_rel}({_wt:+.2f})--> [{_tgt_cat}] {_tgt_lbl}"
                            )

                        if not _edge_lines:
                            st.info("No edges in the Knowledge Graph yet. Add causal relationships first.")
                        else:
                            _graph_ctx = "\n".join(_edge_lines)
                            _forecast_prompt = (
                                "You are a strategic analyst for AGCO/Fendt. "
                                "Below are verified causal relationships from our PESTEL Knowledge Graph. "
                                "Using ONLY these existing edges — no invented sources, no fabricated URLs, "
                                "no external data — identify the 3 most likely NEXT-LEVEL cross-PESTEL "
                                "effects: where does one dimension's shift logically cascade into another?\n\n"
                                "Format each prediction as:\n"
                                "  [SOURCE_DIM] Event → [TARGET_DIM] Predicted effect "
                                "(confidence: High/Medium/Low, reasoning: one sentence)\n\n"
                                "No URLs. No fabricated company names or figures.\n\n"
                                f"EXISTING GRAPH EDGES:\n{_graph_ctx}"
                            )
                            with st.spinner("Generating cross-PESTEL forecast..."):
                                _forecast_text = _call_gemini(_forecast_model, _forecast_prompt, max_tokens=500)
                            if _forecast_text:
                                st.markdown(
                                    f"<div style='background:rgba(150,78,163,0.08);"
                                    f"border-left:3px solid #984ea3;border-radius:6px;"
                                    f"padding:14px;white-space:pre-line;color:#e0e0e0;'>"
                                    f"{_forecast_text}</div>",
                                    unsafe_allow_html=True,
                                )
                                st.caption("⚠️ AI-generated from existing graph edges only. No external sources.")
                            else:
                                _api_unavailable_state("Cross-PESTEL Forecast")
                    except Exception as _fe:
                        st.error(f"Forecast generation failed: {str(_fe)[:200]}")

        except ImportError as e:
            st.error(f"⚠️ Required libraries not installed: {str(e)}")
            st.info("Run: `pip install pyvis networkx`")
        except Exception as e:
            st.error(f"Error loading Knowledge Graph: {str(e)}")

# ===========================
# TAB 6: STRATEGIC REPORTS
# ===========================

with tab6:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0, 204, 255, 0.1); border-left: 4px solid #00ccff; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00ccff; margin: 0;'>📄 Strategic Output Reports</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>C-Suite Markdown Briefs & Board Presentations</p>
    </div>
    """, unsafe_allow_html=True)

    st.caption("**The Writer agent generates executive-ready strategic briefs in Markdown format. These reports synthesize disruption signals into actionable R&D recommendations.**")

    # Search multiple report directories to avoid data-island issues
    report_dirs = [Path('./outputs/reports'), Path('./q2_solution/outputs/reports')]
    all_report_files = []
    for d in report_dirs:
        if d.exists():
            all_report_files.extend(list(d.glob('*.md')))

    if not all_report_files:
        _empty_state(
            "📄", "No Strategic Reports Generated Yet",
            "Reports are automatically produced by the Writer agent during pipeline execution.",
            f"Run <code>python sentinel.py --run-once</code> to generate."
        )
    else:
        report_files = sorted(all_report_files, key=lambda x: x.stat().st_mtime, reverse=True)

        report_names = [f.name for f in report_files]
        selected_report = st.selectbox(
            "📋 Select Report to View",
            options=report_names,
            help="Reports are sorted by modification time (newest first)"
        )

        # Find the actual file object by name
        selected_file = next(f for f in report_files if f.name == selected_report)

        col1, col2 = st.columns(2)
        with col1:
            st.caption(f"**File:** {selected_report}")
        with col2:
            mod_time = datetime.fromtimestamp(selected_file.stat().st_mtime)
            st.caption(f"**Last Modified:** {mod_time.strftime('%Y-%m-%d %H:%M:%S')}")

        st.markdown("---")

        try:
            with open(selected_file, 'r', encoding='utf-8') as f:
                report_content = f.read()

            st.markdown(report_content, unsafe_allow_html=True)

            st.markdown("---")
            st.markdown("**Download report as:**")
            dl_col1, dl_col2, dl_col3 = st.columns(3)

            with dl_col1:
                st.download_button(
                    label="📄 Markdown (.md)",
                    data=report_content,
                    file_name=selected_report,
                    mime="text/markdown"
                )

            with dl_col2:
                # Generate PDF using reportlab
                try:
                    from reportlab.lib.pagesizes import A4
                    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                    from reportlab.lib.units import cm
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
                    from reportlab.lib import colors
                    import io
                    import re

                    def _md_to_pdf_bytes(md_text: str, title: str) -> bytes:
                        buffer = io.BytesIO()
                        doc = SimpleDocTemplate(
                            buffer, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm
                        )
                        styles = getSampleStyleSheet()
                        story = []

                        h1_style = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=16, spaceAfter=10, textColor=colors.HexColor('#003366'))
                        h2_style = ParagraphStyle('H2', parent=styles['Heading2'], fontSize=13, spaceAfter=8, textColor=colors.HexColor('#004499'))
                        h3_style = ParagraphStyle('H3', parent=styles['Heading3'], fontSize=11, spaceAfter=6, textColor=colors.HexColor('#0055AA'))
                        body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=10, spaceAfter=6, leading=14)
                        bullet_style = ParagraphStyle('Bullet', parent=styles['Normal'], fontSize=10, leftIndent=20, spaceAfter=4, leading=13)

                        for line in md_text.split('\n'):
                            line_stripped = line.strip()
                            if not line_stripped:
                                story.append(Spacer(1, 6))
                            elif line_stripped.startswith('### '):
                                txt = line_stripped[4:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h3_style))
                            elif line_stripped.startswith('## '):
                                txt = line_stripped[3:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h2_style))
                            elif line_stripped.startswith('# '):
                                txt = line_stripped[2:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h1_style))
                            elif line_stripped.startswith('---'):
                                story.append(HRFlowable(width="100%", thickness=0.5, color=colors.grey))
                                story.append(Spacer(1, 4))
                            elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                                txt = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', line_stripped[2:])
                                txt = re.sub(r'\*(.*?)\*', r'<i>\1</i>', txt)
                                story.append(Paragraph(f'• {txt}', bullet_style))
                            else:
                                txt = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', line_stripped)
                                txt = re.sub(r'\*(.*?)\*', r'<i>\1</i>', txt)
                                txt = re.sub(r'`(.*?)`', r'<font face="Courier">\1</font>', txt)
                                story.append(Paragraph(txt, body_style))

                        doc.build(story)
                        return buffer.getvalue()

                    pdf_bytes = _md_to_pdf_bytes(report_content, selected_report)
                    pdf_filename = selected_report.replace('.md', '.pdf')
                    st.download_button(
                        label="📕 PDF (.pdf)",
                        data=pdf_bytes,
                        file_name=pdf_filename,
                        mime="application/pdf"
                    )
                except ImportError:
                    st.caption("reportlab not installed — PDF unavailable")
                except Exception as pdf_err:
                    st.caption(f"PDF error: {pdf_err}")

            with dl_col3:
                # Generate PPTX using python-pptx
                try:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    import io
                    import re

                    def _md_to_pptx_bytes(md_text: str, title: str) -> bytes:
                        prs = Presentation()
                        prs.slide_width = Inches(13.33)
                        prs.slide_height = Inches(7.5)

                        blank_layout = prs.slide_layouts[6]
                        title_layout = prs.slide_layouts[0]

                        # Title slide
                        title_slide = prs.slides.add_slide(title_layout)
                        title_slide.shapes.title.text = title.replace('.md', '').replace('_', ' ').title()
                        subtitle_ph = title_slide.placeholders[1]
                        subtitle_ph.text = f"Fendt PESTEL-EL Sentinel — {datetime.now().strftime('%B %d, %Y')}"

                        current_slide = None
                        current_tf = None
                        FENDT_GREEN = RGBColor(0x00, 0x66, 0x33)
                        DARK_BLUE = RGBColor(0x00, 0x33, 0x66)

                        for line in md_text.split('\n'):
                            line_stripped = line.strip()
                            if line_stripped.startswith('# ') or line_stripped.startswith('## '):
                                # New slide per major heading
                                current_slide = prs.slides.add_slide(blank_layout)
                                txBox = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
                                tf = txBox.text_frame
                                tf.word_wrap = True
                                p = tf.paragraphs[0]
                                p.text = line_stripped.lstrip('#').strip()
                                p.runs[0].font.size = Pt(28)
                                p.runs[0].font.bold = True
                                p.runs[0].font.color.rgb = FENDT_GREEN

                                # Body text box for this slide
                                body_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
                                current_tf = body_box.text_frame
                                current_tf.word_wrap = True
                            elif current_tf and line_stripped and not line_stripped.startswith('---'):
                                clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line_stripped)
                                clean = re.sub(r'\*(.*?)\*', r'\1', clean)
                                clean = re.sub(r'`(.*?)`', r'\1', clean)
                                if line_stripped.startswith('- ') or line_stripped.startswith('* '):
                                    clean = '• ' + clean[2:]
                                p = current_tf.add_paragraph()
                                p.text = clean
                                p.font.size = Pt(13)
                                p.space_after = Pt(4)

                        buffer = io.BytesIO()
                        prs.save(buffer)
                        return buffer.getvalue()

                    pptx_bytes = _md_to_pptx_bytes(report_content, selected_report)
                    pptx_filename = selected_report.replace('.md', '.pptx')
                    st.download_button(
                        label="📊 PowerPoint (.pptx)",
                        data=pptx_bytes,
                        file_name=pptx_filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except ImportError:
                    st.caption("python-pptx not installed — PPTX unavailable")
                except Exception as pptx_err:
                    st.caption(f"PPTX error: {pptx_err}")

            # ── Additional export formats: XLS and DOC ──────────────────
            dl_col4, dl_col5 = st.columns(2)

            with dl_col4:
                # Excel export — parse markdown into rows, one per line
                try:
                    import openpyxl
                    import io as _io
                    import re as _re

                    def _md_to_xlsx_bytes(md_text: str, sheet_title: str) -> bytes:
                        wb = openpyxl.Workbook()
                        ws = wb.active
                        ws.title = sheet_title[:31]  # Excel sheet name limit
                        from openpyxl.styles import Font, PatternFill, Alignment
                        HEADER_FILL = PatternFill("solid", fgColor="003366")
                        HEADER_FONT = Font(color="FFFFFF", bold=True, size=12)
                        H2_FONT = Font(color="004499", bold=True, size=11)
                        H3_FONT = Font(color="0055AA", bold=True, size=10)

                        row_num = 1
                        for line in md_text.split('\n'):
                            ls = line.strip()
                            if not ls:
                                row_num += 1
                                continue
                            clean = _re.sub(r'\*\*(.*?)\*\*', r'\1', ls)
                            clean = _re.sub(r'\*(.*?)\*', r'\1', clean)
                            clean = _re.sub(r'`(.*?)`', r'\1', clean)
                            clean = _re.sub(r'^---+$', '─' * 40, clean)
                            cell = ws.cell(row=row_num, column=1, value=clean.lstrip('#- *').strip())
                            cell.alignment = Alignment(wrap_text=True)
                            if ls.startswith('# '):
                                cell.font = HEADER_FONT
                                cell.fill = HEADER_FILL
                            elif ls.startswith('## '):
                                cell.font = H2_FONT
                            elif ls.startswith('### '):
                                cell.font = H3_FONT
                            row_num += 1

                        ws.column_dimensions['A'].width = 120
                        buf = _io.BytesIO()
                        wb.save(buf)
                        return buf.getvalue()

                    xlsx_bytes = _md_to_xlsx_bytes(report_content, selected_report.replace('.md', ''))
                    st.download_button(
                        label="📊 Excel (.xlsx)",
                        data=xlsx_bytes,
                        file_name=selected_report.replace('.md', '.xlsx'),
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    )
                except ImportError:
                    st.download_button(
                        label="📊 Excel (CSV fallback — openpyxl not installed)",
                        data=report_content,
                        file_name=selected_report.replace('.md', '.csv'),
                        mime="text/csv",
                    )
                except Exception as xlsx_err:
                    st.caption(f"Excel error: {xlsx_err}")

            with dl_col5:
                # Word DOC export using python-docx
                try:
                    from docx import Document
                    from docx.shared import Pt, RGBColor as DocxRGB
                    import io as _io
                    import re as _re

                    def _md_to_docx_bytes(md_text: str, doc_title: str) -> bytes:
                        doc = Document()
                        doc.core_properties.title = doc_title
                        doc.core_properties.author = "Fendt PESTEL-EL Sentinel"
                        # Title
                        doc.add_heading(doc_title.replace('.md', '').replace('_', ' ').title(), 0)
                        doc.add_paragraph(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')} UTC")
                        doc.add_paragraph()
                        for line in md_text.split('\n'):
                            ls = line.strip()
                            if not ls:
                                doc.add_paragraph()
                                continue
                            clean = _re.sub(r'\*\*(.*?)\*\*', r'\1', ls)
                            clean = _re.sub(r'\*(.*?)\*', r'\1', clean)
                            clean = _re.sub(r'`(.*?)`', r'\1', clean)
                            if ls.startswith('# '):
                                doc.add_heading(clean.lstrip('# '), level=1)
                            elif ls.startswith('## '):
                                doc.add_heading(clean.lstrip('# '), level=2)
                            elif ls.startswith('### '):
                                doc.add_heading(clean.lstrip('# '), level=3)
                            elif ls.startswith('---'):
                                doc.add_paragraph('─' * 50)
                            elif ls.startswith('- ') or ls.startswith('* '):
                                doc.add_paragraph(clean.lstrip('- '), style='List Bullet')
                            else:
                                doc.add_paragraph(clean)
                        buf = _io.BytesIO()
                        doc.save(buf)
                        return buf.getvalue()

                    docx_bytes = _md_to_docx_bytes(report_content, selected_report)
                    st.download_button(
                        label="📝 Word (.docx)",
                        data=docx_bytes,
                        file_name=selected_report.replace('.md', '.docx'),
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    )
                except ImportError:
                    st.download_button(
                        label="📝 Word (TXT fallback — python-docx not installed)",
                        data=report_content,
                        file_name=selected_report.replace('.md', '.txt'),
                        mime="text/plain",
                    )
                except Exception as docx_err:
                    st.caption(f"Word error: {docx_err}")

            # ── Role-Based AI Summary ─────────────────────────────────────
            st.markdown("---")
            st.markdown("### 🎯 Role-Based Summary")
            st.caption(
                "Select your department. The AI generates a tailored summary of findings, "
                "recommended actions, and 12-month forecasts — grounded strictly in the "
                "signals currently in the database. No external sources or invented facts."
            )

            _DEPT_LENSES = {
                "📢 Marketing": (
                    "competitive positioning, brand differentiation, customer messaging, pricing pressure, "
                    "market share implications"
                ),
                "💼 Sales": (
                    "product portfolio viability, key customer segment risks, revenue impact, "
                    "competitor threat assessment, deal pipeline implications"
                ),
                "🔧 Product Management": (
                    "R&D roadmap prioritization, technology readiness levels, feature gaps vs. competitors, "
                    "patent landscape, build-vs-buy decisions"
                ),
                "🚚 Supply Chain": (
                    "raw material cost volatility, supplier concentration risk, logistics disruption, "
                    "electrification supply chain gaps, component availability for new platforms"
                ),
            }

            _dept_choice = st.selectbox(
                "Select department",
                options=list(_DEPT_LENSES.keys()),
                key="dept_summary_select",
            )

            if st.button("📋 Generate Department Summary", key="dept_summary_btn"):
                _dept_model = _get_gemini_model()
                if not _dept_model:
                    _api_unavailable_state("Department Summary")
                else:
                    # Build signal context from DB (avoid hallucination)
                    _all_sigs = load_all_signals()
                    if _all_sigs:
                        _sig_lines = []
                        for _s in sorted(
                            _all_sigs,
                            key=lambda x: x.get('impact_score') or 0,
                            reverse=True,
                        )[:15]:
                            _sig_lines.append(
                                f"  [{_s.get('primary_dimension','?')}] "
                                f"{_s['title']} — "
                                f"Impact {_s.get('impact_score',0):.2f}, "
                                f"Severity {_s.get('disruption_classification','?')}"
                            )
                        _sig_ctx = "\n".join(_sig_lines)
                    else:
                        _sig_ctx = "(no signals in database)"

                    _dept_prompt = (
                        f"You are a strategic advisor briefing the {_dept_choice.split(' ', 1)[1]} "
                        f"team at AGCO/Fendt. "
                        f"Using ONLY the disruption signals listed below — no invented data, "
                        f"no fabricated URLs, no external sources — produce a structured brief with:\n"
                        f"1. KEY FINDINGS relevant to {_dept_choice.split(' ', 1)[1]} "
                        f"(focus on: {_DEPT_LENSES[_dept_choice]})\n"
                        f"2. RECOMMENDED ACTIONS (3 concrete steps this team should take now)\n"
                        f"3. 12-MONTH FORECAST for this department based strictly on the data\n\n"
                        f"Rules: No URLs. No company names not present in the data. "
                        f"Max 400 words total.\n\n"
                        f"SIGNALS:\n{_sig_ctx}"
                    )

                    with st.spinner(f"Generating {_dept_choice} summary..."):
                        _dept_text = _call_gemini(_dept_model, _dept_prompt, max_tokens=500)

                    if _dept_text:
                        st.markdown(
                            f"<div style='background:rgba(0,204,255,0.07);border-left:3px solid #00ccff;"
                            f"border-radius:6px;padding:16px;white-space:pre-line;color:#e0e0e0;'>"
                            f"<b style='color:#00ccff;'>{_dept_choice} Brief</b><br><br>"
                            f"{_dept_text}</div>",
                            unsafe_allow_html=True,
                        )
                        st.caption("⚠️ AI-generated from database signals only. Always verify before distribution.")
                        # Offer download of this summary
                        st.download_button(
                            label="💾 Download Summary (.txt)",
                            data=f"{_dept_choice} Brief\n{'='*60}\n\n{_dept_text}",
                            file_name=f"dept_summary_{_dept_choice.split()[1].lower()}_{datetime.now().strftime('%Y%m%d')}.txt",
                            mime="text/plain",
                        )
                    else:
                        _api_unavailable_state("Department Summary")

        except Exception as e:
            st.error(f"Error reading report: {str(e)}")

# ===========================
# TAB 7: E-TRACTOR INTELLIGENCE (USE CASE 1)
# ===========================

SEVERITY_COLORS = {
    'CRITICAL': '#ff0066',
    'HIGH': '#ff9933',
    'MODERATE': '#ffff00',
    'LOW': '#00ff88',
}
PRIORITY_BADGE = {
    'CRITICAL': '🔴',
    'HIGH': '🟠',
    'MODERATE': '🟡',
    'LOW': '🟢',
}

with tab7:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0,140,255,0.1); border-left: 4px solid #00ccff; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00ccff; margin: 0;'>📐 Strategic Intelligence Lens</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Deep-dive any strategic topic — powered by semantic search over the live signal database</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0; font-size: 14px;'>
        Select a preset strategic lens or define your own topic. The engine performs semantic search across
        the live intelligence database — no synthetic data, no hardcoded knowledge bases.
        Every result links back to a verified source signal.
        </p>
    </div>
    """, unsafe_allow_html=True)

    if not SEMANTIC_SEARCH_AVAILABLE:
        _empty_state(
            "📐", "Semantic Search Unavailable",
            "semantic_search.py not found in q2_solution/.",
            "Ensure the file exists and scikit-learn is installed."
        )
    elif not signals:
        _empty_state(
            "📡", "No Signals in Database",
            "Run the daily intelligence sweep first to populate the database.",
            "python sentinel.py --run-once"
        )
    else:
        # ── Build search index ───────────────────────────────────────────
        if 'lens_searcher' not in st.session_state:
            _s = SemanticSearch()
            _s.build_index(signals)
            st.session_state['lens_searcher'] = _s
        searcher = st.session_state['lens_searcher']

        import plotly.graph_objects as go

        # ── Preset lenses + custom input ────────────────────────────────
        PRESET_LENSES = {
            "🔋 Electric Tractor Adoption": "electric tractor battery subsidy incentive adoption charging infrastructure",
            "🌿 CAP Reform & Green Deal": "CAP reform EU Green Deal sustainability subsidy reallocation biodiversity",
            "🤖 Precision Farming & Automation": "precision farming autonomous tractor AI robotics GPS sensor automation",
            "🏭 Competitor Intelligence": "John Deere CNH Claas Kubota AGCO competitor product launch patent",
            "📜 Regulatory Compliance": "EU regulation directive compliance mandate emission standard pesticide restriction",
            "💶 Farm Profitability & Subsidies": "farm profitability income subsidy commodity price input cost margin",
            "🌡️ Climate & Weather Risk": "drought flood extreme weather climate change crop yield impact",
            "🔬 AgTech R&D & Innovation": "innovation R&D investment startup funding agtech digital technology patent",
            "✏️ Custom Topic…": "__custom__",
        }

        lens_col, top_k_col = st.columns([4, 1])
        with lens_col:
            selected_lens = st.selectbox(
                "🔎 Select Strategic Lens",
                options=list(PRESET_LENSES.keys()),
                help="Choose a preset strategic topic or define your own custom search query."
            )
        with top_k_col:
            top_k = st.number_input("Max results", min_value=3, max_value=30, value=10, step=1)

        # Custom query input
        if PRESET_LENSES[selected_lens] == "__custom__":
            custom_query = st.text_input(
                "✏️ Enter custom topic or question",
                placeholder="e.g. hydrogen fuel cell tractor EU pilot programs 2026–2028",
                help="Free-form query — the engine finds signals semantically similar to your text."
            )
            active_query = custom_query
        else:
            active_query = PRESET_LENSES[selected_lens]
            st.caption(f"**Search query:** _{active_query}_")

        if not active_query:
            st.info("Enter a custom topic above to begin your analysis.")
        else:
            # ── Run semantic search ─────────────────────────────────────
            with st.spinner("Running semantic search across intelligence database..."):
                results = searcher.search(active_query, top_k=top_k)

            if not results:
                _empty_state(
                    "🔍", "No Matching Signals",
                    "No signals in the current database match this topic.",
                    "Run the intelligence sweep to add more signals: python sentinel.py --run-once"
                )
            else:
                matched_sigs = [sig for sig, _ in results]
                match_df = pd.DataFrame(matched_sigs)

                # ── Summary metrics row ──────────────────────────────────
                n_critical = sum(1 for s in matched_sigs if s.get('disruption_classification') == 'CRITICAL')
                n_high = sum(1 for s in matched_sigs if s.get('disruption_classification') == 'HIGH')
                avg_impact = sum(s.get('impact_score') or 0 for s in matched_sigs) / len(matched_sigs)
                top_dim = match_df['primary_dimension'].value_counts().idxmax() if 'primary_dimension' in match_df.columns else 'N/A'

                m1, m2, m3, m4 = st.columns(4)
                with m1:
                    st.metric("Signals Found", len(results), help="Total signals matching this strategic topic.")
                with m2:
                    st.metric("🔴 Critical / High", f"{n_critical} / {n_high}", help="Signals classified as CRITICAL or HIGH severity within this topic.")
                with m3:
                    st.metric("Avg Impact", f"{avg_impact:.2f}", help="Average Impact score across matched signals.")
                with m4:
                    st.metric("Primary Dimension", top_dim, help="PESTEL dimension with the most signals for this topic.")

                st.markdown("---")

                # ── Two-column layout: chart + signal list ──────────────
                chart_col, list_col = st.columns([1, 2])

                with chart_col:
                    # PESTEL dimension breakdown bar chart
                    if 'primary_dimension' in match_df.columns:
                        dim_counts = match_df['primary_dimension'].value_counts().reset_index()
                        dim_counts.columns = ['Dimension', 'Count']
                        fig_bar = go.Figure(go.Bar(
                            x=dim_counts['Count'],
                            y=dim_counts['Dimension'],
                            orientation='h',
                            marker_color=[PESTEL_COLORS.get(d, '#999') for d in dim_counts['Dimension']],
                            hovertemplate='<b>%{y}</b>: %{x} signals<extra></extra>'
                        ))
                        fig_bar.update_layout(
                            template='plotly_dark', paper_bgcolor='rgba(0,0,0,0)',
                            plot_bgcolor='rgba(0,0,0,0)', height=300,
                            margin=dict(l=0, r=10, t=10, b=0),
                            xaxis=dict(gridcolor='rgba(255,255,255,0.1)', title='Signal Count'),
                            yaxis=dict(gridcolor='rgba(255,255,255,0.05)'),
                            title=dict(text='PESTEL Distribution', font=dict(color='#aaa', size=13), x=0)
                        )
                        st.plotly_chart(fig_bar, use_container_width=True)

                    # Severity breakdown donut
                    sev_counts = match_df['disruption_classification'].value_counts() if 'disruption_classification' in match_df.columns else pd.Series()
                    if len(sev_counts):
                        sev_colors_map = {'CRITICAL': '#ff4b4b', 'HIGH': '#ff9933', 'MODERATE': '#ffdd00', 'LOW': '#00ff88'}
                        fig_pie = go.Figure(go.Pie(
                            labels=sev_counts.index.tolist(),
                            values=sev_counts.values.tolist(),
                            hole=0.6,
                            marker_colors=[sev_colors_map.get(s, '#888') for s in sev_counts.index],
                            hovertemplate='<b>%{label}</b>: %{value} signals<extra></extra>',
                            textinfo='percent+label',
                        ))
                        fig_pie.update_layout(
                            template='plotly_dark', paper_bgcolor='rgba(0,0,0,0)',
                            plot_bgcolor='rgba(0,0,0,0)', height=240,
                            margin=dict(l=0, r=0, t=30, b=0),
                            showlegend=False,
                            title=dict(text='Severity Mix', font=dict(color='#aaa', size=13), x=0)
                        )
                        st.plotly_chart(fig_pie, use_container_width=True)

                with list_col:
                    st.markdown(f"**Top {len(results)} Signals — ranked by semantic relevance to this lens**")

                    for sig, score in results:
                        severity = sig.get('disruption_classification', 'LOW')
                        badge = PRIORITY_BADGE.get(severity, '⚪')
                        sev_color = SEVERITY_COLORS.get(severity, '#888')
                        dim = sig.get('primary_dimension', '')
                        dim_color = PESTEL_COLORS.get(dim, '#999')

                        with st.expander(
                            f"{badge} {sig['title'][:90]}{'…' if len(sig['title']) > 90 else ''}"
                        ):
                            # Header row: dimension badge + scores
                            h1, h2, h3, h4 = st.columns([2, 1, 1, 1])
                            with h1:
                                st.markdown(
                                    f"<span style='background:{dim_color}22;color:{dim_color};"
                                    f"padding:2px 8px;border-radius:4px;font-size:11px;"
                                    f"font-weight:bold;'>{dim}</span>",
                                    unsafe_allow_html=True
                                )
                            with h2:
                                st.metric("Relevance", f"{score:.3f}", help="Semantic similarity score (0–1).")
                            with h3:
                                st.metric("Impact", f"{sig.get('impact_score') or 0:.2f}")
                            with h4:
                                st.metric("Novelty", f"{sig.get('novelty_score') or 0:.2f}")

                            # Content excerpt
                            content = sig.get('content', '')
                            st.caption(content[:500] + '…' if len(content) > 500 else content)

                            # Source link
                            if sig.get('url'):
                                st.markdown(f"[🔗 Open Source]({sig['url']})")

                            # Date
                            if sig.get('date_ingested'):
                                st.caption(f"Ingested: {sig['date_ingested']}")

                st.markdown("---")

                # ── Full signal table export ────────────────────────────
                with st.expander("📋 Full Signal Table (all matched signals)"):
                    table_rows = [
                        {
                            'Title': s['title'],
                            'Dimension': s.get('primary_dimension', ''),
                            'Severity': s.get('disruption_classification', ''),
                            'Impact': round(s.get('impact_score') or 0, 2),
                            'Novelty': round(s.get('novelty_score') or 0, 2),
                            'Velocity': round(s.get('velocity_score') or 0, 2),
                            'Relevance Score': round(score, 3),
                            'Source URL': s.get('url', ''),
                            'Date': s.get('date_ingested', ''),
                        }
                        for s, score in results
                    ]
                    st.dataframe(
                        pd.DataFrame(table_rows),
                        use_container_width=True,
                        hide_index=True,
                        column_config={
                            "Title": st.column_config.TextColumn("Signal", width="large"),
                            "Impact": st.column_config.NumberColumn("Impact", format="%.2f"),
                            "Novelty": st.column_config.NumberColumn("Novelty", format="%.2f"),
                            "Velocity": st.column_config.NumberColumn("Velocity", format="%.2f"),
                            "Relevance Score": st.column_config.NumberColumn("Relevance", format="%.3f"),
                            "Source URL": st.column_config.LinkColumn("🔗 Source"),
                        }
                    )

# ===========================
# FOOTER
# ===========================

st.sidebar.markdown("---")
st.sidebar.caption("© 2026 Fendt Strategic Intelligence")
st.sidebar.caption("Powered by Claude Code & SQLite")
