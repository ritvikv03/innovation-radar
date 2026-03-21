"""
Fendt PESTEL-EL Sentinel Dashboard
====================================

Modern Streamlit dashboard powered by SQLite backend (q2_solution).

**Phase 4 Enterprise Features:**
- BLUF AI Executive Narrative (or rule-based fallback)
- Conversational AI Strategy Sparring (or prototype fallback)

Features 6 tabs:
1. Executive Summary: AI/rule-based BLUF + metrics + critical disruptions
2. Innovation Radar: Interactive Plotly visualization with dimension filters
3. Live Signal Feed: Searchable dataframe of all signals
4. The Inquisition: Conversational AI strategy dialogue (or prototype mode)
5. Knowledge Graph: Causal interdependencies (Table View + Network Graph)
6. Strategic Reports: C-suite markdown briefs viewer

**Prototype Mode:** Works without GEMINI_API_KEY using intelligent fallbacks.
**Production Mode:** Configure API key for full Claude AI capabilities.
"""

import streamlit as st
import pandas as pd
from datetime import datetime
from pathlib import Path
import sys
import os
from typing import List, Dict
from dotenv import load_dotenv

# Graceful Gemini import — app must render even if package is absent
try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

load_dotenv()

# Add q2_solution to path — multiple strategies for local vs Streamlit Cloud
PROJECT_ROOT = Path(__file__).parent.absolute()
_q2_path = str(PROJECT_ROOT / "q2_solution")
if _q2_path not in sys.path:
    sys.path.insert(0, _q2_path)

try:
    from database import SignalDatabase
except ImportError:
    sys.path.insert(0, str(PROJECT_ROOT))
    from q2_solution.database import SignalDatabase

try:
    from use_case_1_etractor import ETractorIntelligence
    from semantic_search import SemanticSearch
    ETRACTOR_AVAILABLE = True
except ImportError:
    try:
        from q2_solution.use_case_1_etractor import ETractorIntelligence
        from q2_solution.semantic_search import SemanticSearch
        ETRACTOR_AVAILABLE = True
    except ImportError:
        ETRACTOR_AVAILABLE = False

try:
    from innovation_radar import InnovationRadar, PESTEL_COLORS
except ImportError:
    try:
        from q2_solution.innovation_radar import InnovationRadar, PESTEL_COLORS
    except ImportError:
        from q2_solution.innovation_radar import InnovationRadar
        PESTEL_COLORS = {
            'POLITICAL': '#e41a1c', 'ECONOMIC': '#377eb8', 'SOCIAL': '#4daf4a',
            'TECHNOLOGICAL': '#984ea3', 'ENVIRONMENTAL': '#ff7f00', 'LEGAL': '#ffff33',
            'INNOVATION': '#00ccff', 'SOCIAL_MEDIA': '#ff69b4',
        }

# Preferred Gemini models in priority order (newest stable first)
_GEMINI_MODELS = ["gemini-2.0-flash", "gemini-1.5-flash", "gemini-1.5-flash-8b"]


def get_api_key() -> str:
    """Get Gemini API key from Streamlit secrets (cloud) or env (local)."""
    try:
        if "GEMINI_API_KEY" in st.secrets:
            return st.secrets["GEMINI_API_KEY"]
    except Exception:
        pass
    return os.getenv("GEMINI_API_KEY", "")


@st.cache_data(ttl=1800)  # Re-check every 30 minutes
def check_gemini_health() -> tuple:
    """
    Probe the Gemini API once and cache the result for 30 minutes.

    Returns:
        (is_healthy: bool, model_name: str, status_msg: str)
    """
    api_key = get_api_key()
    if not api_key or not GEMINI_AVAILABLE:
        return False, "", "No API key configured"

    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        for model_name in _GEMINI_MODELS:
            try:
                model = genai.GenerativeModel(model_name)
                resp = model.generate_content(
                    "Reply with exactly two words: STATUS OK",
                    generation_config={"max_output_tokens": 10},
                )
                if resp.text:
                    return True, model_name, f"OK ({model_name})"
            except Exception as e:
                err = str(e)
                if "429" in err or "quota" in err.lower() or "RESOURCE_EXHAUSTED" in err:
                    return False, "", "Daily quota exhausted — rule-based fallback active"
                # Other error: try next model
                continue
        return False, "", "All Gemini models unavailable"
    except Exception as e:
        return False, "", f"Gemini error: {str(e)[:80]}"


def _get_gemini_model():
    """Return a ready GenerativeModel using the healthiest available model, or None."""
    healthy, model_name, _ = check_gemini_health()
    if not healthy:
        return None
    import google.generativeai as genai
    genai.configure(api_key=get_api_key())
    return genai.GenerativeModel(model_name)


def _empty_state(icon: str, title: str, message: str, action: str = "") -> None:
    """Render a stylized, enterprise-grade empty state card."""
    action_html = (
        f"<p style='color:#888;margin:10px 0 0 0;font-size:13px;'>{action}</p>"
        if action else ""
    )
    st.markdown(f"""
    <div style='padding:48px 24px;text-align:center;background:rgba(255,255,255,0.03);
                border:1px dashed rgba(255,255,255,0.12);border-radius:12px;margin:16px 0;'>
        <div style='font-size:52px;margin-bottom:14px;'>{icon}</div>
        <h3 style='color:#cccccc;margin:0 0 8px 0;font-weight:600;'>{title}</h3>
        <p style='color:#888888;margin:0;font-size:15px;'>{message}</p>
        {action_html}
    </div>
    """, unsafe_allow_html=True)


def _api_unavailable_state(feature: str) -> None:
    """Render a stylized warning when Gemini API is unavailable (missing key or quota)."""
    _, _, status_msg = check_gemini_health()
    is_quota = "quota" in status_msg.lower() or "exhausted" in status_msg.lower()
    icon = "⏳" if is_quota else "⚙️"
    detail = (
        "Daily free-tier quota exhausted. The rule-based fallback is active. "
        "The quota resets at midnight UTC, or upgrade to a paid Gemini plan for uninterrupted AI."
        if is_quota else
        f"<strong>{feature}</strong> requires a <code>GEMINI_API_KEY</code> in Streamlit secrets "
        f"or your <code>.env</code> file. A rule-based fallback is active in the meantime."
    )
    st.markdown(f"""
    <div style='padding:16px 20px;background:rgba(255,153,0,0.08);
                border-left:4px solid #ff9900;border-radius:8px;margin:12px 0;'>
        <h4 style='color:#ff9900;margin:0 0 4px 0;'>{icon} AI Fallback Mode — {feature}</h4>
        <p style='color:#ccc;margin:0;font-size:13px;'>{detail}</p>
        <p style='color:#888;margin:4px 0 0 0;font-size:12px;'>Status: {status_msg}</p>
    </div>
    """, unsafe_allow_html=True)

# ===========================
# PAGE CONFIGURATION
# ===========================
st.set_page_config(
    page_title="Fendt PESTEL-EL Sentinel",
    page_icon="🛰️",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for premium dark theme
st.markdown("""
<style>
    .main {
        background-color: #0e1117;
    }
    .stApp {
        color: #e0e0e0;
    }
    .metric-card {
        padding: 20px;
        border-radius: 10px;
        background: rgba(255, 255, 255, 0.05);
        border: 1px solid rgba(255, 255, 255, 0.1);
        margin-bottom: 15px;
    }
    .critical-alert {
        padding: 15px;
        border-left: 4px solid #ff4b4b;
        background: rgba(255, 75, 75, 0.1);
        border-radius: 5px;
        margin-bottom: 10px;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
    }
    .stTabs [data-baseweb="tab"] {
        padding: 10px 20px;
        background-color: rgba(255, 255, 255, 0.05);
        border-radius: 5px;
    }
</style>
""", unsafe_allow_html=True)

# ===========================
# DATABASE CONNECTION
# ===========================

@st.cache_resource
def get_db():
    """Get cached database connection."""
    return SignalDatabase()

@st.cache_data(ttl=60)  # Cache for 1 minute
def load_all_signals():
    """Load all signals from database."""
    db = get_db()
    return db.get_all_signals()

@st.cache_data(ttl=60)
def get_db_stats():
    """Get database statistics."""
    db = get_db()
    return db.get_database_stats()

# ===========================
# BLUF AI EXECUTIVE NARRATIVE
# ===========================

def generate_bluf_narrative(db_stats: Dict, signals: List[Dict]) -> str:
    """
    Generate Bottom Line Up Front (BLUF) executive narrative.

    Uses Gemini API if available, otherwise generates intelligent rule-based summary.

    Analyzes database statistics and critical signals to produce a 3-sentence
    hard-hitting strategic summary for C-suite consumption.

    Args:
        db_stats: Database statistics from get_database_stats()
        signals: List of all signals from database

    Returns:
        str: 3-sentence BLUF narrative
    """
    api_key = get_api_key()

    # Analyze signal distribution
    df = pd.DataFrame(signals)

    # Count by dimension
    dimension_counts = df['primary_dimension'].value_counts().to_dict()
    top_dimension = max(dimension_counts, key=dimension_counts.get)
    top_dimension_pct = (dimension_counts[top_dimension] / len(df)) * 100

    # Count by severity
    severity_counts = df['disruption_classification'].value_counts().to_dict()
    critical_count = severity_counts.get('CRITICAL', 0)
    high_count = severity_counts.get('HIGH', 0)

    # Calculate average scores
    avg_velocity = df['velocity_score'].mean() if 'velocity_score' in df.columns else 0

    # Get top critical signal
    critical_df = df[df['disruption_classification'] == 'CRITICAL'].sort_values('impact_score', ascending=False)
    top_threat = critical_df.iloc[0]['title'] if len(critical_df) > 0 else None

    # Use rule-based fallback if Gemini is unavailable or quota exhausted
    gemini_model = _get_gemini_model()
    if not gemini_model:
        # Sentence 1: Overall status and volume
        sentence1 = f"Currently tracking {db_stats['total_signals']} disruption signals with {critical_count} CRITICAL and {high_count} HIGH priority threats demanding immediate attention."

        # Sentence 2: Dimension analysis
        if top_dimension_pct > 40:
            sentence2 = f"Signals are heavily skewed toward {top_dimension} dimension ({top_dimension_pct:.0f}% of total), indicating concentrated regulatory/market pressure in this area."
        else:
            sentence2 = f"Disruptions are distributed across PESTEL dimensions with {top_dimension} showing highest concentration ({top_dimension_pct:.0f}%)."

        # Sentence 3: Velocity and top threat
        if avg_velocity > 0.6:
            velocity_desc = "rapidly accelerating"
        elif avg_velocity > 0.4:
            velocity_desc = "building momentum"
        else:
            velocity_desc = "emerging gradually"

        if top_threat:
            sentence3 = f"Velocity indicators show {velocity_desc} disruption patterns, with most critical threat identified as: {top_threat[:80]}..."
        else:
            sentence3 = f"Velocity indicators show {velocity_desc} disruption patterns across the intelligence landscape."

        return f"{sentence1} {sentence2} {sentence3}"

    # AI path — Gemini is healthy, generate sophisticated narrative
    avg_impact = df['impact_score'].mean() if 'impact_score' in df.columns else 0
    avg_novelty = df['novelty_score'].mean() if 'novelty_score' in df.columns else 0

    context = f"""DATABASE STATISTICS:
- Total Signals: {db_stats['total_signals']}
- Date Range: {db_stats['date_range']['earliest']} to {db_stats['date_range']['latest']}

PESTEL DISTRIBUTION:
{chr(10).join(f'- {dim}: {count} signals ({count/db_stats["total_signals"]*100:.1f}%)' for dim, count in dimension_counts.items())}

SEVERITY BREAKDOWN:
{chr(10).join(f'- {severity}: {count}' for severity, count in severity_counts.items())}

AVERAGE SCORES: Impact {avg_impact:.2f} | Velocity {avg_velocity:.2f} | Novelty {avg_novelty:.2f}

TOP CRITICAL THREAT: {top_threat}"""

    prompt = (
        "You are a strategic advisor delivering a BLUF briefing to AGCO/Fendt's senior leadership. "
        "Analyze the data below and write EXACTLY 3 hard-hitting, data-driven sentences. "
        "No preamble, no titles, no bullet points — ONLY the 3 sentences.\n\n"
        f"{context}"
    )

    try:
        response = gemini_model.generate_content(
            prompt, generation_config={"max_output_tokens": 300}
        )
        return response.text.strip()
    except Exception:
        # Fall through to rule-based on any error
        pass

    # Fallback if AI call fails mid-flight
    return (
        f"Currently tracking {db_stats['total_signals']} disruption signals with "
        f"{critical_count} CRITICAL and {high_count} HIGH priority threats. "
        f"{top_dimension} dimension shows highest concentration ({top_dimension_pct:.0f}%). "
        f"{'Priority threat: ' + top_threat[:80] + '...' if top_threat else 'No critical threats detected.'}"
    )

# ===========================
# THE INQUISITION - GEMINI API
# ===========================

def generate_strategic_questions(critical_signals: List[Dict]) -> List[str]:
    """
    Generate 3-5 hard-hitting strategic questions for Fendt C-suite
    using Gemini API based on critical signals.

    Args:
        critical_signals: List of critical disruption signals

    Returns:
        List of strategic questions
    """
    _fallback_questions = [
        "If EU electrification mandates accelerate to 2027, which Fendt R&D programs face obsolescence?",
        "Which competitor is closest to commercializing autonomous precision farming, and what is Fendt's countermove?",
        "How does the CAP subsidy reallocation toward sustainability affect Fendt's core diesel tractor revenue?",
    ]

    gemini_model = _get_gemini_model()
    if not gemini_model:
        return _fallback_questions

    signal_summaries = "\n".join(
        f"- {sig['title']} [{sig['primary_dimension']}] "
        f"(Disruption: {sig.get('disruption_classification', 'N/A')}, "
        f"Impact: {sig.get('impact_score', 0):.2f})"
        for sig in critical_signals[:10]
    )

    prompt = (
        "You are a strategic advisor to AGCO/Fendt's senior leadership. "
        "Based on these disruption signals, generate 3-5 hard-hitting strategic questions "
        "numbered 1-5. No preamble. ONLY the questions.\n\n"
        f"SIGNALS:\n{signal_summaries}"
    )

    try:
        response = gemini_model.generate_content(
            prompt, generation_config={"max_output_tokens": 400}
        )
        questions = []
        for line in response.text.strip().split('\n'):
            line = line.strip()
            if line and (line[0].isdigit() or line.startswith('-')):
                q = line.lstrip('0123456789.-) *').strip()
                if q:
                    questions.append(q)
        return questions if questions else _fallback_questions
    except Exception:
        return _fallback_questions

# ===========================
# MAIN DASHBOARD
# ===========================

st.markdown("""
<div style='text-align: center; padding: 20px; background: linear-gradient(135deg, rgba(0,100,0,0.1) 0%, rgba(0,50,0,0.2) 100%); border-radius: 10px; margin-bottom: 20px;'>
    <h1 style='color: #00ff88; margin: 0;'>🛰️ Fendt PESTEL-EL Sentinel</h1>
    <h3 style='color: #00ccff; margin: 10px 0;'>Strategic Intelligence War Room</h3>
    <p style='color: #999; margin: 0;'>Autonomous Agricultural Disruption Detection for AGCO/Fendt C-Suite</p>
</div>
""", unsafe_allow_html=True)

# Sidebar
st.sidebar.title("System Status")
stats = get_db_stats()

if stats['total_signals'] > 0:
    st.sidebar.metric("Total Signals", stats['total_signals'])

    if stats['date_range']:
        st.sidebar.caption(f"Earliest: {stats['date_range']['earliest']}")
        st.sidebar.caption(f"Latest: {stats['date_range']['latest']}")
else:
    st.sidebar.warning("No signals in database yet.")
    st.sidebar.info("Run the daily intelligence sweep to populate data.")

st.sidebar.markdown("---")

# Gemini API health status
_g_healthy, _g_model, _g_status = check_gemini_health()
_g_color = "#00ff88" if _g_healthy else "#ff9933"
_g_icon = "🟢" if _g_healthy else "🟡"
st.sidebar.markdown(
    f"<span style='font-size:12px;color:{_g_color};'>{_g_icon} AI: {_g_status}</span>",
    unsafe_allow_html=True,
)

st.sidebar.markdown("---")
st.sidebar.markdown("**PESTEL-EL Breakdown**")
if stats['signals_per_dimension']:
    for dim, count in stats['signals_per_dimension'].items():
        color = PESTEL_COLORS.get(dim, '#999999')
        st.sidebar.markdown(
            f"<span style='color:{color};font-size:16px;'>●</span> "
            f"<span style='font-size:13px;'><b>{dim}</b>: {count}</span>",
            unsafe_allow_html=True
        )

# Load signals
signals = load_all_signals()

# ===========================
# TAB LAYOUT
# ===========================

tab1, tab2, tab3, tab4, tab5, tab6, tab7 = st.tabs([
    "📊 Intelligence Overview",
    "🎯 Innovation Radar",
    "📡 PESTEL Signal Monitor",
    "⚔️ The Inquisition",
    "🕸️ Knowledge Graph",
    "📄 Strategic Reports",
    "🔋 E-Tractor Intelligence",
])

# ===========================
# TAB 1: EXECUTIVE SUMMARY
# ===========================

with tab1:
    st.subheader("Intelligence Overview")
    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0; font-size: 14px;'>
        This tab provides an automated overview of disruption signals detected, classified, and scored by the Sentinel.
        Signals are ingested from live EU agricultural intelligence sources (EUR-Lex, Eurostat, news APIs) and evaluated
        across eight PESTEL-EL dimensions. Metrics below reflect the current contents of the intelligence database —
        they are updated each time the daily sweep runs.
        </p>
    </div>
    """, unsafe_allow_html=True)

    if not signals:
        _empty_state(
            "📡", "Intelligence Database Empty",
            "No disruption signals have been ingested yet.",
            "Run the daily intelligence sweep to populate the radar: <code>python sentinel.py --run-once</code>"
        )
    else:
        # ===========================
        # BLUF AI EXECUTIVE NARRATIVE
        # ===========================

        st.markdown("""
        <div style='padding: 20px; background: linear-gradient(135deg, rgba(0,100,200,0.15) 0%, rgba(0,50,100,0.2) 100%); border-left: 4px solid #00ccff; border-radius: 10px; margin-bottom: 10px;'>
            <h3 style='color: #00ccff; margin: 0 0 6px 0;'>🎯 AI-Generated Signal Synthesis (BLUF)</h3>
            <p style='color: #888; margin: 0; font-size: 13px;'>
                Bottom Line Up Front — a 3-sentence strategic summary of the signals currently in the intelligence database,
                generated automatically from PESTEL distribution, severity counts, and disruption scores.
                <b>Note:</b> "Novelty" measures how unique each signal is <em>relative to previously ingested signals in this database</em>,
                not whether the event is globally new. A high novelty score means the Sentinel has not seen this type of signal before.
            </p>
        </div>
        """, unsafe_allow_html=True)

        # Generate or retrieve BLUF from session state
        if 'bluf_narrative' not in st.session_state or st.button("🔄 Regenerate AI Narrative", help="Click to generate a fresh AI analysis of current intelligence data"):
            with st.spinner("Generating AI executive narrative..."):
                st.session_state.bluf_narrative = generate_bluf_narrative(stats, signals)

        # Display BLUF with help tooltip
        st.markdown(
            f"""
            <div style='font-size: 18px; line-height: 1.8; color: #e0e0e0; padding: 15px; background: rgba(0,0,0,0.3); border-radius: 8px; margin-bottom: 25px;'>
                {st.session_state.bluf_narrative}
            </div>
            """,
            unsafe_allow_html=True
        )

        if not get_api_key() or not GEMINI_AVAILABLE:
            _api_unavailable_state("BLUF AI Executive Narrative")
        else:
            st.caption("ℹ️ **AI-Generated Synthesis:** Dynamically generated by Gemini analyzing PESTEL distribution, severity breakdown, and mathematical disruption scores.")

        st.markdown("---")

        # Convert to DataFrame for easier manipulation
        df = pd.DataFrame(signals)

        # Top-level metrics
        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric(
                "Total Signals", len(signals),
                help="Total number of classified disruption signals ingested from monitored EU agricultural intelligence sources (EUR-Lex, Eurostat, news APIs). Each signal is a detected event or trend that has been scored for Impact, Novelty, and Velocity."
            )

        with col2:
            critical_count = len(df[df['disruption_classification'] == 'CRITICAL'])
            st.metric(
                "🔴 Critical", critical_count,
                help="CRITICAL signals are classified as immediate-horizon threats (0–12 months) with a composite disruption score above the CRITICAL threshold. These require urgent review by senior leadership (CEO, CFO, CTO, VP R&D)."
            )

        with col3:
            high_count = len(df[df['disruption_classification'] == 'HIGH'])
            st.metric(
                "🟠 High", high_count,
                help="HIGH signals are near-term threats (12–24 months) with elevated disruption scores. These require strategic planning and resource allocation decisions at the VP/Director level."
            )

        with col4:
            # Average disruption score
            if 'impact_score' in df.columns and df['impact_score'].notna().any():
                avg_impact = df['impact_score'].mean()
                st.metric("Avg Impact", f"{avg_impact:.2f}", help="Average Impact across all ingested signals. This tracks the overall magnitude of industry disruption currently affecting the European agricultural landscape.")
            else:
                st.metric("Avg Impact", "N/A", help="Average Impact across all ingested signals. This tracks the overall magnitude of industry disruption currently affecting the European agricultural landscape.")

        st.markdown("---")

        # Time-Series Velocity Chart
        st.markdown("### 📈 Temporal Velocity Tracking")
        st.markdown("""
        <div style='padding: 10px 14px; background: rgba(255,255,255,0.03); border-radius: 6px; margin-bottom: 10px; border-left: 3px solid rgba(0,204,255,0.4);'>
            <p style='color: #aaa; margin: 0; font-size: 13px;'>
            <b>What this chart shows:</b> The number of new disruption signals ingested per day over time.
            A rising curve indicates accelerating activity in the intelligence feed (more events detected).
            A flat or declining curve indicates a quieter period. This view enables tracking whether disruption
            momentum is building, stable, or decelerating — distinct from any individual signal's Velocity score
            (which measures a single signal's acceleration relative to historical averages).
            </p>
        </div>
        """, unsafe_allow_html=True)

        # Group signals by date and count
        if 'date_ingested' in df.columns:
            df['date_parsed'] = pd.to_datetime(df['date_ingested']).dt.date
            daily_counts = df.groupby('date_parsed').size().reset_index(name='count')
            daily_counts = daily_counts.sort_values('date_parsed')
            
            # Format explicitly as strings to prevent Plotly from rendering continuous sub-second ticks
            daily_counts['date_str'] = daily_counts['date_parsed'].apply(lambda x: x.strftime('%b %d, %Y'))

            if len(daily_counts) > 0:
                # Create Plotly line chart for better aesthetics
                import plotly.graph_objects as go

                fig = go.Figure()
                fig.add_trace(go.Scatter(
                    x=daily_counts['date_str'],
                    y=daily_counts['count'],
                    mode='lines+markers',
                    name='Daily Signals',
                    line=dict(color='#00ff88', width=3),
                    marker=dict(size=8, color='#00ff88', line=dict(color='white', width=1)),
                    hovertemplate='<b>%{x}</b><br>Signals: %{y}<extra></extra>'
                ))

                fig.update_layout(
                    template='plotly_dark',
                    plot_bgcolor='rgba(0,0,0,0)',
                    paper_bgcolor='rgba(0,0,0,0)',
                    xaxis_title='Date',
                    yaxis_title='Signal Count',
                    height=250,
                    margin=dict(l=0, r=0, t=10, b=0),
                    xaxis=dict(
                        gridcolor='rgba(255,255,255,0.1)',
                        type='category',
                        tickangle=-45
                    ),
                    yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
                    showlegend=False
                )

                st.plotly_chart(fig, use_container_width=True)
            else:
                st.info("Insufficient data for temporal chart. Need signals from multiple days.")
        else:
            st.warning("date_ingested field not found in signals.")

        st.markdown("---")

        # Critical Disruptions Section
        st.markdown("### 🚨 Top 3 Critical Disruptions")
        st.markdown("""
        <div style='padding: 10px 14px; background: rgba(255,75,75,0.06); border-radius: 6px; margin-bottom: 10px; border-left: 3px solid rgba(255,75,75,0.5);'>
            <p style='color: #aaa; margin: 0; font-size: 13px;'>
            Signals classified as <b>CRITICAL</b> (composite disruption score in the top tier, immediate 0–12 month horizon).
            Ranked by Impact score. Each card shows the signal's content excerpt, source link, and individual
            Impact / Novelty / Velocity sub-scores. These are the three signals most likely to require
            a strategic response from senior leadership in the near term.
            </p>
        </div>
        """, unsafe_allow_html=True)

        critical_signals = df[df['disruption_classification'] == 'CRITICAL'].to_dict('records')

        if critical_signals:
            # Sort by impact score (descending)
            critical_signals.sort(
                key=lambda x: x.get('impact_score', 0) if x.get('impact_score') is not None else 0,
                reverse=True
            )

            # Show top 3 with detailed view
            for i, signal in enumerate(critical_signals[:3]):  # Show top 3
                with st.expander(
                    f"🔴 {signal['title']} [{signal['primary_dimension']}]",
                    expanded=(i < 3)  # Expand first 3
                ):
                    col_a, col_b = st.columns([2, 1])

                    with col_a:
                        st.markdown(f"**Source:** [{signal['source']}]({signal['url']})")
                        st.markdown(f"**Date:** {signal['date_ingested']}")
                        st.caption(signal['content'][:500] + "..." if len(signal['content']) > 500 else signal['content'])

                    with col_b:
                        st.metric(
                            "Impact",
                            f"{signal.get('impact_score', 0):.2f}",
                            help="Impact (0.0-1.0) measures magnitude. Calculated algorithmically based on cross-PESTEL reach (how many dimensions the signal affects) and high-leverage triggers like legal mandates or tech breakthroughs. Crucial for flagging existential industry threats."
                        )
                        st.metric(
                            "Novelty",
                            f"{signal.get('novelty_score', 0):.2f}",
                            help="Novelty (0.0-1.0) measures uniqueness. Calculated via inverse-similarity text matching against historical signals in the SQLite database. Crucial for separating true emerging anomalies from repetitive background noise."
                        )
                        st.metric(
                            "Velocity",
                            f"{signal.get('velocity_score', 0):.2f}",
                            help="Velocity (0.0-1.0) measures acceleration. Derived mathematically by comparing recent 30-day signal volume against the trailing 6-month historical average. Crucial to prove a trend is actively gaining momentum, rather than just being a one-off event."
                        )

            # Show remaining critical signals count
            if len(critical_signals) > 3:
                st.info(f"ℹ️ **{len(critical_signals) - 3} additional critical signals** detected. View all in the Live Signal Feed tab.")
        else:
            st.success("✅ No critical disruptions detected. System is monitoring for emerging threats.")

# ===========================
# TAB 2: INNOVATION RADAR
# ===========================

with tab2:
    st.subheader("Innovation Radar - Industry Disruption Map")

    if not signals:
        _empty_state(
            "🎯", "Radar Awaiting Data",
            "No signals available to plot on the disruption radar.",
            "Run the daily intelligence sweep to populate the radar: <code>python sentinel.py --run-once</code>"
        )
    else:
        st.markdown("""
        **Time Horizons:**
        - 🔴 **12 Month**: Immediate action required
        - 🟡 **24 Month**: Pilot and trial phase
        - 🟢 **36 Month**: Assess and monitor

        **Sectors** (8 PESTEL-EL pillars): Political · Economic · Social · Technological · Environmental · Legal · Innovation · Social Media
        """)

        # Dimension filter widget
        df_radar = pd.DataFrame(signals)
        available_dimensions = sorted(df_radar['primary_dimension'].unique().tolist())

        selected_dimensions = st.multiselect(
            "🎛️ Filter by PESTEL Dimension",
            options=available_dimensions,
            default=available_dimensions,
            help="Select which PESTEL dimensions to display on the radar. Deselect to focus on specific areas."
        )

        if not selected_dimensions:
            st.warning("Please select at least one dimension to display the radar.")
        else:
            # Prepare signals for radar
            radar_signals = []
            # Filter signals by selected dimensions
            filtered_signals = [sig for sig in signals if sig['primary_dimension'] in selected_dimensions]

            for sig in filtered_signals:
                # Map disruption classification to time horizon
                classification = sig.get('disruption_classification', 'LOW')
                if classification == 'CRITICAL':
                    horizon = '12_MONTH'
                elif classification == 'HIGH':
                    horizon = '24_MONTH'
                else:
                    horizon = '36_MONTH'

                # Calculate composite disruption score
                impact = sig.get('impact_score', 0) if sig.get('impact_score') is not None else 0
                novelty = sig.get('novelty_score', 0) if sig.get('novelty_score') is not None else 0
                velocity = sig.get('velocity_score', 0) if sig.get('velocity_score') is not None else 0

                disruption_score = (impact * 0.5 + novelty * 0.3 + velocity * 0.2)

                radar_signals.append({
                    'title': sig['title'],
                    'primary_dimension': sig['primary_dimension'],
                    'time_horizon': horizon,
                    'disruption_score': disruption_score,
                    'classification': classification
                })

            # Generate radar
            try:
                radar = InnovationRadar()
                fig = radar.create_radar(radar_signals)
                st.plotly_chart(fig, use_container_width=True)
            except Exception as e:
                st.error(f"Error generating radar: {str(e)}")
                st.info("Ensure plotly is installed: pip install plotly")

            # Signal table with source links
            st.markdown("---")
            st.markdown("### 📋 Plotted Signals — Source Reference")
            st.caption("All signals currently displayed on the radar, with direct links to source documents.")

            # Build table from original signals (with URL) for filtered dimensions
            radar_table_data = []
            for sig in filtered_signals:
                classification = sig.get('disruption_classification', 'LOW')
                if classification == 'CRITICAL':
                    horizon_label = '🔴 12 Month'
                elif classification == 'HIGH':
                    horizon_label = '🟡 24 Month'
                else:
                    horizon_label = '🟢 36 Month'

                impact = sig.get('impact_score', 0) if sig.get('impact_score') is not None else 0
                novelty = sig.get('novelty_score', 0) if sig.get('novelty_score') is not None else 0
                velocity = sig.get('velocity_score', 0) if sig.get('velocity_score') is not None else 0
                disruption_score = round(impact * 0.5 + novelty * 0.3 + velocity * 0.2, 3)

                radar_table_data.append({
                    'Title': sig['title'],
                    'Dimension': sig['primary_dimension'],
                    'Severity': classification,
                    'Time Horizon': horizon_label,
                    'Disruption Score': disruption_score,
                    'Source URL': sig.get('url', ''),
                })

            radar_table_df = pd.DataFrame(radar_table_data).sort_values('Disruption Score', ascending=False)
            st.dataframe(
                radar_table_df,
                use_container_width=True,
                hide_index=True,
                column_config={
                    "Title": st.column_config.TextColumn("Signal Title", width="large"),
                    "Dimension": st.column_config.TextColumn("PESTEL Dimension", width="small"),
                    "Severity": st.column_config.TextColumn("Severity", width="small"),
                    "Time Horizon": st.column_config.TextColumn("Time Horizon", width="small",
                        help="Horizon is derived from severity: CRITICAL → 12 months, HIGH → 24 months, MODERATE/LOW → 36 months."),
                    "Disruption Score": st.column_config.NumberColumn("Score", format="%.3f", width="small",
                        help="Composite score = Impact×0.5 + Novelty×0.3 + Velocity×0.2"),
                    "Source URL": st.column_config.LinkColumn("🔗 Source", width="medium",
                        help="Direct link to the original source document."),
                }
            )

# ===========================
# TAB 3: LIVE SIGNAL FEED
# ===========================

with tab3:
    st.subheader("PESTEL Signal Monitor")
    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0 0 8px 0; font-size: 14px;'>
        Raw intelligence signals ingested from monitored sources, fully searchable and filterable.
        Each signal has been classified into a PESTEL-EL dimension and scored on three axes:
        </p>
        <ul style='color: #888; margin: 0; font-size: 13px; padding-left: 18px;'>
            <li><b>Impact (0–1):</b> Magnitude of the disruption — how many PESTEL dimensions are affected and whether it involves a high-leverage trigger (legal mandate, technology breakthrough, subsidy change).</li>
            <li><b>Novelty (0–1):</b> Uniqueness relative to signals already in this database — calculated via text similarity matching against historical entries. High novelty = the Sentinel has not seen this type of event before.</li>
            <li><b>Velocity (0–1):</b> Momentum — mathematical comparison of 30-day signal volume vs. the trailing 6-month average. High velocity = this topic is actively accelerating.</li>
        </ul>
        <p style='color: #888; margin: 8px 0 0 0; font-size: 13px;'>
        <b>Severity classification</b> is derived from the composite disruption score
        (Impact×0.5 + Novelty×0.3 + Velocity×0.2): CRITICAL ≥ 0.7 · HIGH ≥ 0.5 · MODERATE ≥ 0.3 · LOW &lt; 0.3.
        Duplicate signals (same source URL) are filtered at display time.
        </p>
    </div>
    """, unsafe_allow_html=True)

    if not signals:
        _empty_state(
            "📡", "Signal Feed Empty",
            "No disruption signals have been ingested yet.",
            "Run the daily intelligence sweep: <code>python sentinel.py --run-once</code>"
        )
    else:
        df = pd.DataFrame(signals)

        # Deduplicate by URL at display time (keep first / most recent occurrence)
        if 'url' in df.columns:
            df_deduped = df.drop_duplicates(subset=['url'], keep='first').copy()
        else:
            df_deduped = df.copy()

        # Search and filter controls
        col_search, col_filter = st.columns([2, 1])

        with col_search:
            search_term = st.text_input("🔍 Search signals (title, content, source)", "")

        with col_filter:
            dimension_filter = st.multiselect(
                "Filter by dimension",
                options=df_deduped['primary_dimension'].unique().tolist(),
                default=[]
            )

        # Apply filters
        filtered_df = df_deduped.copy()

        if search_term:
            filtered_df = filtered_df[
                filtered_df['title'].str.contains(search_term, case=False, na=False) |
                filtered_df['content'].str.contains(search_term, case=False, na=False) |
                filtered_df['source'].str.contains(search_term, case=False, na=False)
            ]

        if dimension_filter:
            filtered_df = filtered_df[filtered_df['primary_dimension'].isin(dimension_filter)]

        duplicates_removed = len(df) - len(df_deduped)
        caption_parts = [f"Showing {len(filtered_df)} of {len(df_deduped)} signals"]
        if duplicates_removed > 0:
            caption_parts.append(f"({duplicates_removed} duplicate URLs removed)")
        st.caption(" ".join(caption_parts))

        # Display configuration — include URL as clickable link
        display_columns = [
            'title', 'primary_dimension', 'disruption_classification',
            'impact_score', 'novelty_score', 'velocity_score',
            'date_ingested', 'source', 'url'
        ]

        # Only show columns that exist
        display_columns = [col for col in display_columns if col in filtered_df.columns]

        # Display dataframe
        st.dataframe(
            filtered_df[display_columns].sort_values('date_ingested', ascending=False),
            use_container_width=True,
            hide_index=True,
            column_config={
                "title": st.column_config.TextColumn(
                    "Title", width="large",
                    help="Signal title as extracted from the source document."
                ),
                "primary_dimension": st.column_config.TextColumn(
                    "Dimension", width="small",
                    help="The primary PESTEL-EL dimension this signal belongs to, determined by the Classifier agent based on signal content and keyword triggers."
                ),
                "disruption_classification": st.column_config.TextColumn(
                    "Severity", width="small",
                    help="Severity classification derived from composite disruption score: CRITICAL (≥0.7), HIGH (≥0.5), MODERATE (≥0.3), LOW (<0.3). Reflects urgency of response, not absolute danger."
                ),
                "impact_score": st.column_config.NumberColumn(
                    "Impact",
                    format="%.2f",
                    help="Impact (0.0–1.0): Magnitude of disruption. Calculated based on cross-PESTEL reach (how many dimensions the signal touches) and high-leverage triggers like legal mandates or technology breakthroughs."
                ),
                "novelty_score": st.column_config.NumberColumn(
                    "Novelty",
                    format="%.2f",
                    help="Novelty (0.0–1.0): Uniqueness vs. existing database entries. Uses text similarity matching — high novelty means the Sentinel has not previously seen this type of signal. Does NOT indicate global newness."
                ),
                "velocity_score": st.column_config.NumberColumn(
                    "Velocity",
                    format="%.2f",
                    help="Velocity (0.0–1.0): Momentum score. Compares 30-day signal count for this topic vs. 6-month historical average. High velocity = trend actively accelerating."
                ),
                "date_ingested": st.column_config.TextColumn(
                    "Date Ingested", width="small",
                    help="Date the signal was detected and ingested by the Sentinel pipeline."
                ),
                "source": st.column_config.TextColumn(
                    "Source", width="small",
                    help="Publication or data provider name (e.g., EUR-Lex, Eurostat, AgriPulse)."
                ),
                "url": st.column_config.LinkColumn(
                    "🔗 Source Link", width="medium",
                    help="Direct link to the original source document. Click to open in browser."
                ),
            }
        )

        # Export option
        st.download_button(
            label="📥 Download as CSV",
            data=filtered_df.to_csv(index=False),
            file_name=f"fendt_signals_{datetime.now().strftime('%Y%m%d')}.csv",
            mime="text/csv"
        )

# ===========================
# TAB 4: THE INQUISITION (CONVERSATIONAL AI)
# ===========================

with tab4:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(255, 75, 75, 0.1); border-left: 4px solid #ff4b4b; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #ff4b4b; margin: 0;'>⚔️ The Inquisition</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Conversational AI Strategy Sparring for C-Suite</p>
    </div>
    """, unsafe_allow_html=True)
    st.caption("**Engage in strategic dialogue with Claude. Ask follow-up questions, challenge assumptions, and explore alternative scenarios based on current disruption intelligence.**")

    if not signals:
        _empty_state(
            "⚔️", "The Inquisition Awaits Intelligence",
            "No signals available. The Inquisition requires live disruption data to interrogate.",
            "Run the daily intelligence sweep first: <code>python sentinel.py --run-once</code>"
        )
    else:
        df = pd.DataFrame(signals)
        # Accept both CRITICAL and HIGH signals for analysis
        high_priority_signals = df[df['disruption_classification'].isin(['CRITICAL', 'HIGH'])].to_dict('records')

        if not high_priority_signals:
            st.warning("No high-priority signals detected. The Inquisition requires at least one CRITICAL or HIGH signal.")
            st.info("Run the daily intelligence sweep to populate data.")
        else:
            critical_count = len(df[df['disruption_classification'] == 'CRITICAL'])
            high_count = len(df[df['disruption_classification'] == 'HIGH'])
            st.markdown(f"**Analyzing {critical_count} CRITICAL and {high_count} HIGH signals ({len(high_priority_signals)} total)...**")

            # Initialize conversation state
            if 'inquisition_messages' not in st.session_state:
                st.session_state.inquisition_messages = []

            if 'inquisition_initialized' not in st.session_state:
                st.session_state.inquisition_initialized = False

            # Show awaiting state on first load — only fire on explicit button click
            if not st.session_state.inquisition_initialized:
                _empty_state(
                    "🔮", "Ready to Interrogate",
                    "Click the button above to generate hard-hitting strategic questions from your live intelligence data.",
                    "The Inquisition analyzes your CRITICAL and HIGH signals and forces C-suite reflection."
                )

            if st.button("🔮 Start New Strategic Session", type="primary"):
                with st.spinner("Consulting Claude for strategic insights..."):
                    questions = generate_strategic_questions(high_priority_signals)

                # Reset conversation
                st.session_state.inquisition_messages = []

                # Add system context
                context_summary = f"I've analyzed {len(high_priority_signals)} high-priority disruption signals ({critical_count} CRITICAL, {high_count} HIGH) across the PESTEL framework for Fendt/AGCO."

                # Format initial message with questions
                initial_message = f"{context_summary}\n\nHere are the critical strategic questions I need you to confront:\n\n"
                for i, question in enumerate(questions, 1):
                    initial_message += f"**{i}. {question}**\n\n"

                initial_message += "How would you like to respond? You can:\n- Answer a specific question\n- Challenge my assumptions\n- Ask for clarification\n- Explore alternative scenarios"

                # Add to conversation history
                st.session_state.inquisition_messages.append({
                    "role": "assistant",
                    "content": initial_message
                })

                st.session_state.inquisition_initialized = True
                st.rerun()

            # Display conversation history
            st.markdown("### 💬 Strategic Dialogue")

            for message in st.session_state.inquisition_messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            # Chat input
            if prompt := st.chat_input("Type your response or question (e.g., 'Assume we can't change diesel roadmap until 2028. How does that affect Question 2?')"):
                # Add user message to history
                st.session_state.inquisition_messages.append({
                    "role": "user",
                    "content": prompt
                })

                # Display user message
                with st.chat_message("user"):
                    st.markdown(prompt)

                # Generate AI response
                with st.chat_message("assistant"):
                    with st.spinner("Generating strategic analysis..."):
                        inq_model = _get_gemini_model()
                        if not inq_model:
                            _, _, status = check_gemini_health()
                            st.warning(
                                f"AI unavailable ({status}). The Inquisition requires Gemini to respond. "
                                "Please try again after midnight UTC when the quota resets, "
                                "or configure a paid Gemini API key in Streamlit secrets."
                            )
                        else:
                            try:
                                signal_context = "\n".join(
                                    f"- {sig['title']} [{sig['primary_dimension']}] "
                                    f"(Impact: {sig.get('impact_score', 0):.2f})"
                                    for sig in high_priority_signals[:10]
                                )
                                chat_context = (
                                    "You are a strategic advisor to AGCO/Fendt's senior leadership "
                                    "in a strategic dialogue about agricultural industry disruptions.\n\n"
                                    f"INTELLIGENCE CONTEXT:\n{signal_context}\n\n"
                                )
                                for msg in st.session_state.inquisition_messages[:-1]:
                                    chat_context += f"{msg['role'].upper()}: {msg['content']}\n\n"
                                chat_context += f"USER: {prompt}\n\nRespond strategically:"

                                response = inq_model.generate_content(
                                    chat_context,
                                    generation_config={"max_output_tokens": 600},
                                )
                                full_response = response.text
                                st.markdown(full_response)
                                st.session_state.inquisition_messages.append({
                                    "role": "assistant",
                                    "content": full_response,
                                })
                            except Exception as e:
                                st.error(f"Error: {str(e)[:200]}")

            # Show signal context
            with st.expander("📋 View High-Priority Signals Being Analyzed"):
                for sig in high_priority_signals[:10]:
                    st.markdown(f"**{sig['title']}** [{sig['primary_dimension']}]")
                    st.caption(f"Impact: {sig.get('impact_score', 0):.2f} | Source: {sig['source']}")
                    st.markdown("---")

# ===========================
# TAB 5: KNOWLEDGE GRAPH
# ===========================

with tab5:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0, 255, 136, 0.1); border-left: 4px solid #00ff88; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00ff88; margin: 0;'>🕸️ Knowledge Graph</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Causal Interdependencies Across PESTEL Dimensions</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown("""
    <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 18px; border: 1px solid rgba(255,255,255,0.08);'>
        <p style='color: #aaaaaa; margin: 0 0 8px 0; font-size: 14px;'><b>How it works:</b></p>
        <ul style='color: #888; margin: 0; font-size: 13px; padding-left: 18px;'>
            <li><b>Nodes</b> represent named entities extracted from intelligence signals — a regulatory policy, a market shift, a technology, or a physical event.</li>
            <li><b>Edges</b> represent causal relationships between two entities, identified by the Analyst agent from signal content and verbatim quotes. Each edge has a <b>weight</b> (−1.0 to +1.0): positive = one event accelerates or enables the other; negative = one event restricts or counters the other.</li>
            <li><b>Temporal decay:</b> Edge weights decay with a 90-day half-life — older, unconfirmed relationships automatically lose influence. Edges below 0.05 are pruned.</li>
            <li><b>Provenance:</b> Every edge must have a verified <code>source_url</code> and an <code>exact_quote</code> from the source (minimum 10 characters) — required for EU Data Act 2026 compliance.</li>
            <li><b>Color coding</b> (node color = PESTEL dimension; edge color = relationship strength and direction — see legend below the graph).</li>
        </ul>
    </div>
    """, unsafe_allow_html=True)

    graph_path = Path('./data/graph.json')

    if not graph_path.exists():
        _empty_state(
            "🕸️", "Knowledge Graph Not Yet Initialized",
            "No causal relationships have been mapped between PESTEL entities.",
            "The Graph is populated by the Analyst agents during pipeline execution: <code>python sentinel.py --run-once</code>"
        )
    else:
        try:
            import json
            import networkx as nx

            # Load graph
            with open(graph_path, 'r') as f:
                graph_data = json.load(f)

            # Mirror links/edges to prevent NetworkX version KeyErrors across environments
            if 'links' in graph_data and 'edges' not in graph_data:
                graph_data['edges'] = graph_data['links']
            elif 'edges' in graph_data and 'links' not in graph_data:
                graph_data['links'] = graph_data['edges']

            G = nx.node_link_graph(graph_data)

            # Display stats
            col1, col2, col3 = st.columns(3)
            with col1:
                st.metric("Nodes (Entities)", G.number_of_nodes(), help="Nodes represent distinct entities derived from the intelligence feed. A node can be a regulatory policy, an economic shift, a specific technology, or a physical environment change. The more nodes tracked, the broader our horizon scanning.")
            with col2:
                st.metric("Edges (Relationships)", G.number_of_edges(), help="Edges represent the mathematical links connecting two distinct nodes. This explicitly tracks how one event (e.g., a Political change) ripples out to trigger another event (e.g., an Economic cost).")
            with col3:
                if G.number_of_edges() > 0:
                    avg_weight = sum([attrs.get('weight', 0) for _, _, attrs in G.edges(data=True)]) / G.number_of_edges()
                    st.metric("Avg Edge Weight", f"{avg_weight:.2f}", help="Weight (-1.0 to 1.0) measures the strength of the causal ripple. Positive values accelerate or amplify a trend. Negative values decelerate, regulate, or restrict a trend.")
                else:
                    st.metric("Avg Edge Weight", "N/A", help="Weight (-1.0 to 1.0) measures the strength of the causal ripple. Positive values accelerate or amplify a trend. Negative values decelerate, regulate, or restrict a trend.")

            st.markdown("---")

            # View selection
            view_mode = st.radio(
                "📊 Visualization Mode",
                ["📋 Table View (Recommended)", "🕸️ Network Graph"],
                help="Table View is clearer for large graphs. Network Graph shows spatial relationships."
            )

            if view_mode == "📋 Table View (Recommended)":
                st.markdown("### Causal Relationships")
                st.caption("**Each row shows one causal link. Read as: Source → Relationship → Target**")

                causality_data = []
                for source, target, attrs in G.edges(data=True):
                    source_label = G.nodes[source].get('label', source)
                    source_cat = G.nodes[source].get('category', 'N/A')
                    target_label = G.nodes[target].get('label', target)
                    target_cat = G.nodes[target].get('category', 'N/A')

                    relationship = attrs.get('relationship', 'RELATES_TO')
                    weight = attrs.get('weight', 0.0)
                    quote = attrs.get('exact_quote', 'N/A')
                    source_url = attrs.get('source_url', 'N/A')

                    causality_data.append({
                        'Source': f"{source_label} [{source_cat}]",
                        'Relationship': relationship,
                        'Target': f"{target_label} [{target_cat}]",
                        'Weight': weight,
                        'Effect': '🟢 Positive' if weight > 0 else '🔴 Negative' if weight < 0 else '⚪ Neutral',
                        'Evidence Quote': quote[:100] + '...' if len(quote) > 100 else quote,
                        'Source URL': source_url
                    })

                causality_df = pd.DataFrame(causality_data)
                causality_df['abs_weight'] = causality_df['Weight'].abs()
                causality_df = causality_df.sort_values('abs_weight', ascending=False).drop('abs_weight', axis=1)

                st.dataframe(
                    causality_df,
                    use_container_width=True,
                    hide_index=True,
                    column_config={
                        "Source": st.column_config.TextColumn("📍 Source Node", width="medium", help="The inciting PESTEL entity that triggered the disruption."),
                        "Relationship": st.column_config.TextColumn("⚡ Causal Verb", width="small", help="The physical action explaining how the source affects the target."),
                        "Target": st.column_config.TextColumn("🎯 Target Node", width="medium", help="The resulting PESTEL entity that is being impacted."),
                        "Weight": st.column_config.NumberColumn("📊 Weight", format="%.2f", width="small", help="The mathematical strength (-1.0 to 1.0) of the ripple effect."),
                        "Effect": st.column_config.TextColumn("↗️ Effect Type", width="small", help="Whether the trend is being accelerated (positive) or restricted (negative)."),
                        "Evidence Quote": st.column_config.TextColumn("📜 Quote", width="large", help="Exact text pulled from the intelligence agent justifying this link."),
                        "Source URL": st.column_config.LinkColumn("🔗 Source", width="medium", help="Direct link to the source document for verification.")
                    }
                )

                st.download_button(
                    label="📥 Download Causality Table as CSV",
                    data=causality_df.to_csv(index=False),
                    file_name=f"knowledge_graph_causality_{datetime.now().strftime('%Y%m%d')}.csv",
                    mime="text/csv"
                )

            else:
                from pyvis.network import Network
                import tempfile
                import streamlit.components.v1 as components

                st.markdown("### Interactive Network Visualization")
                st.caption("**Drag nodes to explore. Click for details. Arrows show causality direction.**")

                col_filter1, col_filter2 = st.columns(2)

                with col_filter1:
                    min_weight = st.slider(
                        "Minimum Edge Weight (Strength)",
                        min_value=0.0,
                        max_value=1.0,
                        value=0.1,
                        step=0.1,
                        help="Filter out weak relationships to reduce clutter"
                    )

                with col_filter2:
                    all_categories = set()
                    for _, attrs in G.nodes(data=True):
                        all_categories.add(attrs.get('category', 'UNKNOWN'))

                    selected_categories = st.multiselect(
                        "Filter by PESTEL Dimension",
                        options=sorted(all_categories),
                        default=sorted(all_categories),
                        help="Show only nodes from selected dimensions"
                    )

                net = Network(
                    height='650px',
                    width='100%',
                    bgcolor='#0e1117',
                    font_color='#ffffff',
                    directed=True
                )

                net.barnes_hut(
                    gravity=-10000,
                    central_gravity=0.1,
                    spring_length=250,
                    spring_strength=0.001,
                    damping=0.5
                )

                category_colors = PESTEL_COLORS  # Single source of truth from innovation_radar.py

                visible_nodes = set()
                for node_id, attrs in G.nodes(data=True):
                    category = attrs.get('category', 'UNKNOWN')

                    if category not in selected_categories:
                        continue

                    label = attrs.get('label', node_id)
                    color = category_colors.get(category, '#999999')

                    display_label = label[:30] + '...' if len(label) > 30 else label

                    net.add_node(
                        node_id,
                        label=display_label,
                        title=f"<b>{label}</b><br>{category}",
                        color=color,
                        size=35,
                        font={'size': 18, 'color': '#ffffff', 'face': 'arial', 'bold': True}
                    )
                    visible_nodes.add(node_id)

                for source, target, attrs in G.edges(data=True):
                    if source not in visible_nodes or target not in visible_nodes:
                        continue

                    weight = attrs.get('weight', 0.0)

                    if abs(weight) < min_weight:
                        continue

                    relationship = attrs.get('relationship', 'RELATES_TO')

                    if weight > 0.5:
                        edge_color = '#00ff88'
                    elif weight > 0:
                        edge_color = '#00ccff'
                    elif weight < -0.5:
                        edge_color = '#ff0066'
                    elif weight < 0:
                        edge_color = '#ff9933'
                    else:
                        edge_color = '#999999'

                    net.add_edge(
                        source,
                        target,
                        title=f"{relationship} (weight: {weight:.2f})",
                        color=edge_color,
                        width=abs(weight) * 5,
                        arrows={'to': {'enabled': True, 'scaleFactor': 1.5}}
                    )

                with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.html') as f:
                    net.save_graph(f.name)
                    with open(f.name, 'r') as html_file:
                        html_content = html_file.read()

                    components.html(html_content, height=670, scrolling=False)

                st.markdown("---")
                st.markdown("### 🎨 Edge Color Legend")
                col_l1, col_l2, col_l3, col_l4 = st.columns(4)

                with col_l1:
                    st.markdown("🟢 **Strong Positive** (Weight > 0.5)")
                with col_l2:
                    st.markdown("🔵 **Moderate Positive** (Weight 0 to 0.5)")
                with col_l3:
                    st.markdown("🟠 **Moderate Negative** (Weight -0.5 to 0)")
                with col_l4:
                    st.markdown("🔴 **Strong Negative** (Weight < -0.5)")

        except ImportError as e:
            st.error(f"⚠️ Required libraries not installed: {str(e)}")
            st.info("Run: `pip install pyvis networkx`")
        except Exception as e:
            st.error(f"Error loading Knowledge Graph: {str(e)}")

# ===========================
# TAB 6: STRATEGIC REPORTS
# ===========================

with tab6:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0, 204, 255, 0.1); border-left: 4px solid #00ccff; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00ccff; margin: 0;'>📄 Strategic Output Reports</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>C-Suite Markdown Briefs & Board Presentations</p>
    </div>
    """, unsafe_allow_html=True)

    st.caption("**The Writer agent generates executive-ready strategic briefs in Markdown format. These reports synthesize disruption signals into actionable R&D recommendations.**")

    # Search multiple report directories to avoid data-island issues
    report_dirs = [Path('./outputs/reports'), Path('./q2_solution/outputs/reports')]
    all_report_files = []
    for d in report_dirs:
        if d.exists():
            all_report_files.extend(list(d.glob('*.md')))

    if not all_report_files:
        _empty_state(
            "📄", "No Strategic Reports Generated Yet",
            "Reports are automatically produced by the Writer agent during pipeline execution.",
            f"Run <code>python sentinel.py --run-once</code> to generate."
        )
    else:
        report_files = sorted(all_report_files, key=lambda x: x.stat().st_mtime, reverse=True)

        report_names = [f.name for f in report_files]
        selected_report = st.selectbox(
            "📋 Select Report to View",
            options=report_names,
            help="Reports are sorted by modification time (newest first)"
        )

        # Find the actual file object by name
        selected_file = next(f for f in report_files if f.name == selected_report)

        col1, col2 = st.columns(2)
        with col1:
            st.caption(f"**File:** {selected_report}")
        with col2:
            mod_time = datetime.fromtimestamp(selected_file.stat().st_mtime)
            st.caption(f"**Last Modified:** {mod_time.strftime('%Y-%m-%d %H:%M:%S')}")

        st.markdown("---")

        try:
            with open(selected_file, 'r', encoding='utf-8') as f:
                report_content = f.read()

            st.markdown(report_content, unsafe_allow_html=True)

            st.markdown("---")
            st.markdown("**Download report as:**")
            dl_col1, dl_col2, dl_col3 = st.columns(3)

            with dl_col1:
                st.download_button(
                    label="📄 Markdown (.md)",
                    data=report_content,
                    file_name=selected_report,
                    mime="text/markdown"
                )

            with dl_col2:
                # Generate PDF using reportlab
                try:
                    from reportlab.lib.pagesizes import A4
                    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
                    from reportlab.lib.units import cm
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
                    from reportlab.lib import colors
                    import io
                    import re

                    def _md_to_pdf_bytes(md_text: str, title: str) -> bytes:
                        buffer = io.BytesIO()
                        doc = SimpleDocTemplate(
                            buffer, pagesize=A4,
                            rightMargin=2*cm, leftMargin=2*cm,
                            topMargin=2*cm, bottomMargin=2*cm
                        )
                        styles = getSampleStyleSheet()
                        story = []

                        h1_style = ParagraphStyle('H1', parent=styles['Heading1'], fontSize=16, spaceAfter=10, textColor=colors.HexColor('#003366'))
                        h2_style = ParagraphStyle('H2', parent=styles['Heading2'], fontSize=13, spaceAfter=8, textColor=colors.HexColor('#004499'))
                        h3_style = ParagraphStyle('H3', parent=styles['Heading3'], fontSize=11, spaceAfter=6, textColor=colors.HexColor('#0055AA'))
                        body_style = ParagraphStyle('Body', parent=styles['Normal'], fontSize=10, spaceAfter=6, leading=14)
                        bullet_style = ParagraphStyle('Bullet', parent=styles['Normal'], fontSize=10, leftIndent=20, spaceAfter=4, leading=13)

                        for line in md_text.split('\n'):
                            line_stripped = line.strip()
                            if not line_stripped:
                                story.append(Spacer(1, 6))
                            elif line_stripped.startswith('### '):
                                txt = line_stripped[4:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h3_style))
                            elif line_stripped.startswith('## '):
                                txt = line_stripped[3:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h2_style))
                            elif line_stripped.startswith('# '):
                                txt = line_stripped[2:].replace('**', '').replace('*', '')
                                story.append(Paragraph(txt, h1_style))
                            elif line_stripped.startswith('---'):
                                story.append(HRFlowable(width="100%", thickness=0.5, color=colors.grey))
                                story.append(Spacer(1, 4))
                            elif line_stripped.startswith('- ') or line_stripped.startswith('* '):
                                txt = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', line_stripped[2:])
                                txt = re.sub(r'\*(.*?)\*', r'<i>\1</i>', txt)
                                story.append(Paragraph(f'• {txt}', bullet_style))
                            else:
                                txt = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', line_stripped)
                                txt = re.sub(r'\*(.*?)\*', r'<i>\1</i>', txt)
                                txt = re.sub(r'`(.*?)`', r'<font face="Courier">\1</font>', txt)
                                story.append(Paragraph(txt, body_style))

                        doc.build(story)
                        return buffer.getvalue()

                    pdf_bytes = _md_to_pdf_bytes(report_content, selected_report)
                    pdf_filename = selected_report.replace('.md', '.pdf')
                    st.download_button(
                        label="📕 PDF (.pdf)",
                        data=pdf_bytes,
                        file_name=pdf_filename,
                        mime="application/pdf"
                    )
                except ImportError:
                    st.caption("reportlab not installed — PDF unavailable")
                except Exception as pdf_err:
                    st.caption(f"PDF error: {pdf_err}")

            with dl_col3:
                # Generate PPTX using python-pptx
                try:
                    from pptx import Presentation
                    from pptx.util import Inches, Pt
                    from pptx.dml.color import RGBColor
                    import io
                    import re

                    def _md_to_pptx_bytes(md_text: str, title: str) -> bytes:
                        prs = Presentation()
                        prs.slide_width = Inches(13.33)
                        prs.slide_height = Inches(7.5)

                        blank_layout = prs.slide_layouts[6]
                        title_layout = prs.slide_layouts[0]

                        # Title slide
                        title_slide = prs.slides.add_slide(title_layout)
                        title_slide.shapes.title.text = title.replace('.md', '').replace('_', ' ').title()
                        subtitle_ph = title_slide.placeholders[1]
                        subtitle_ph.text = f"Fendt PESTEL-EL Sentinel — {datetime.now().strftime('%B %d, %Y')}"

                        current_slide = None
                        current_tf = None
                        FENDT_GREEN = RGBColor(0x00, 0x66, 0x33)
                        DARK_BLUE = RGBColor(0x00, 0x33, 0x66)

                        for line in md_text.split('\n'):
                            line_stripped = line.strip()
                            if line_stripped.startswith('# ') or line_stripped.startswith('## '):
                                # New slide per major heading
                                current_slide = prs.slides.add_slide(blank_layout)
                                txBox = current_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1.0))
                                tf = txBox.text_frame
                                tf.word_wrap = True
                                p = tf.paragraphs[0]
                                p.text = line_stripped.lstrip('#').strip()
                                p.runs[0].font.size = Pt(28)
                                p.runs[0].font.bold = True
                                p.runs[0].font.color.rgb = FENDT_GREEN

                                # Body text box for this slide
                                body_box = current_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))
                                current_tf = body_box.text_frame
                                current_tf.word_wrap = True
                            elif current_tf and line_stripped and not line_stripped.startswith('---'):
                                clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line_stripped)
                                clean = re.sub(r'\*(.*?)\*', r'\1', clean)
                                clean = re.sub(r'`(.*?)`', r'\1', clean)
                                if line_stripped.startswith('- ') or line_stripped.startswith('* '):
                                    clean = '• ' + clean[2:]
                                p = current_tf.add_paragraph()
                                p.text = clean
                                p.font.size = Pt(13)
                                p.space_after = Pt(4)

                        buffer = io.BytesIO()
                        prs.save(buffer)
                        return buffer.getvalue()

                    pptx_bytes = _md_to_pptx_bytes(report_content, selected_report)
                    pptx_filename = selected_report.replace('.md', '.pptx')
                    st.download_button(
                        label="📊 PowerPoint (.pptx)",
                        data=pptx_bytes,
                        file_name=pptx_filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                except ImportError:
                    st.caption("python-pptx not installed — PPTX unavailable")
                except Exception as pptx_err:
                    st.caption(f"PPTX error: {pptx_err}")

        except Exception as e:
            st.error(f"Error reading report: {str(e)}")

# ===========================
# TAB 7: E-TRACTOR INTELLIGENCE (USE CASE 1)
# ===========================

SEVERITY_COLORS = {
    'CRITICAL': '#ff0066',
    'HIGH': '#ff9933',
    'MODERATE': '#ffff00',
    'LOW': '#00ff88',
}
PRIORITY_BADGE = {
    'CRITICAL': '🔴',
    'HIGH': '🟠',
    'MODERATE': '🟡',
    'LOW': '🟢',
}

with tab7:
    st.markdown("""
    <div style='text-align: center; padding: 15px; background: rgba(0,180,80,0.1); border-left: 4px solid #00cc55; border-radius: 5px; margin-bottom: 20px;'>
        <h2 style='color: #00cc55; margin: 0;'>🔋 Use Case 1: Electric Tractor Intelligence</h2>
        <p style='color: #ccc; margin: 5px 0 0 0;'>Subsidies · Business Impact · Sales & Marketing Recommendations · Customer Needs</p>
    </div>
    """, unsafe_allow_html=True)

    if not ETRACTOR_AVAILABLE:
        _empty_state(
            "🔋", "E-Tractor Module Unavailable",
            "use_case_1_etractor.py or semantic_search.py not found.",
            "Ensure both files are in the q2_solution/ directory."
        )
    else:
        etractor = ETractorIntelligence()
        searcher = SemanticSearch()
        searcher.build_index(signals)

        # ── Overview metrics ────────────────────────────────────────────
        impact = etractor.generate_business_impact(signals)
        relevant_signals = etractor.get_relevant_signals(signals)
        competitor_signals = etractor.get_competitor_signals(signals)

        risk_color = {'HIGH': '#ff0066', 'MODERATE': '#ff9933', 'LOW': '#00ff88'}.get(
            impact['overall_risk_level'], '#999'
        )

        col_r1, col_r2, col_r3, col_r4 = st.columns(4)
        with col_r1:
            st.metric(
                "🔋 E-Tractor Signals", impact['total_relevant_signals'],
                help="Signals in the database that contain electric tractor, battery, subsidy, or related keywords."
            )
        with col_r2:
            st.metric(
                "🔴 Critical / High", impact['critical_high_count'],
                help="Subset classified as CRITICAL or HIGH severity — requiring immediate strategic attention."
            )
        with col_r3:
            st.metric(
                "🏭 Competitor Signals", len(competitor_signals),
                help="Signals mentioning competitor brands (John Deere, CNH, Claas, Kubota, etc.) in the e-tractor context."
            )
        with col_r4:
            st.metric(
                "⚡ Avg Impact Score", f"{impact['avg_impact_score']:.2f}",
                help="Average Impact score (0–1) across e-tractor-relevant signals."
            )

        st.markdown(
            f"<div style='padding:10px 14px;background:rgba(255,255,255,0.04);border-radius:6px;"
            f"border-left:4px solid {risk_color};margin-bottom:18px;'>"
            f"<b style='color:{risk_color};'>Overall Risk Level: {impact['overall_risk_level']}</b>"
            f"<span style='color:#888;font-size:13px;'> — based on {impact['critical_high_count']} "
            f"CRITICAL/HIGH signals out of {impact['total_relevant_signals']} e-tractor-relevant signals detected.</span>"
            f"</div>",
            unsafe_allow_html=True
        )

        st.markdown("---")

        # ── Seven sections via expanders ────────────────────────────────
        section_tabs = st.tabs([
            "🌍 Subsidy Programs",
            "📊 Business Impact",
            "💼 Sales Recommendations",
            "📣 Marketing Recommendations",
            "🚜 Customer Needs",
            "🔍 Semantic Search",
        ])

        # ── SECTION 1: Subsidy Programs ──────────────────────────────
        with section_tabs[0]:
            st.markdown("### Available Subsidy & Incentive Programs")
            st.caption(
                "Static knowledge base covering current (2024–2027) programs. "
                "Amounts and eligibility criteria are indicative — verify with national agricultural agencies before advising customers."
            )

            country_selected = st.selectbox(
                "🌍 Select Country / Region",
                options=etractor.get_all_countries(),
                help="Select a geography to view applicable subsidy programs for electric tractor adoption."
            )

            country_data = etractor.get_subsidy_programs(country_selected)
            st.markdown(
                f"<div style='padding:10px 14px;background:rgba(0,180,80,0.07);border-radius:6px;"
                f"border-left:3px solid #00cc55;margin-bottom:14px;'>"
                f"<p style='color:#aaa;margin:0;font-size:14px;'>{country_data.get('headline', '')}</p>"
                f"</div>",
                unsafe_allow_html=True
            )

            for prog in country_data.get('programs', []):
                rel_color = SEVERITY_COLORS.get(prog['relevance'], '#888')
                badge = PRIORITY_BADGE.get(prog['relevance'], '⚪')
                with st.expander(f"{badge} {prog['name']} — {prog['type']}"):
                    col_a, col_b = st.columns([3, 1])
                    with col_a:
                        st.markdown(f"**Description:** {prog['description']}")
                        st.markdown(f"**Timeline:** {prog['timeline']}")
                        if prog.get('url'):
                            st.markdown(f"**Source:** [{prog['url']}]({prog['url']})")
                    with col_b:
                        st.markdown(
                            f"<div style='padding:10px;background:rgba(255,255,255,0.04);border-radius:6px;text-align:center;'>"
                            f"<p style='color:#888;margin:0;font-size:11px;'>AMOUNT</p>"
                            f"<p style='color:#e0e0e0;margin:4px 0;font-size:13px;font-weight:bold;'>{prog['amount']}</p>"
                            f"<p style='color:{rel_color};margin:0;font-size:12px;font-weight:bold;'>{prog['relevance']} RELEVANCE</p>"
                            f"</div>",
                            unsafe_allow_html=True
                        )

            if not country_data.get('programs'):
                st.info("No subsidy data available for this selection.")

        # ── SECTION 2: Business Impact ───────────────────────────────
        with section_tabs[1]:
            st.markdown("### Business Impact Assessment — AGCO/Fendt")
            st.caption(
                "Analysis derived from live intelligence signals in the database. "
                "Refreshes automatically as new signals are ingested."
            )

            if not relevant_signals:
                _empty_state(
                    "🔋", "No E-Tractor Signals Detected",
                    "Run the intelligence sweep to populate e-tractor relevant signals.",
                    "python sentinel.py --run-once"
                )
            else:
                # Dimension distribution
                if impact['dimension_distribution']:
                    st.markdown("**📊 Signal Distribution by PESTEL Dimension**")
                    dim_df = pd.DataFrame([
                        {'Dimension': k, 'Signal Count': v}
                        for k, v in sorted(impact['dimension_distribution'].items(), key=lambda x: x[1], reverse=True)
                    ])
                    import plotly.graph_objects as go
                    fig_dim = go.Figure(go.Bar(
                        x=dim_df['Dimension'],
                        y=dim_df['Signal Count'],
                        marker_color=[PESTEL_COLORS.get(d, '#999') for d in dim_df['Dimension']],
                        hovertemplate='<b>%{x}</b>: %{y} signals<extra></extra>'
                    ))
                    fig_dim.update_layout(
                        template='plotly_dark', paper_bgcolor='rgba(0,0,0,0)',
                        plot_bgcolor='rgba(0,0,0,0)', height=280,
                        margin=dict(l=0, r=0, t=10, b=0),
                        yaxis=dict(gridcolor='rgba(255,255,255,0.1)'),
                        xaxis=dict(gridcolor='rgba(255,255,255,0.05)')
                    )
                    st.plotly_chart(fig_dim, use_container_width=True)

                # Top signals table
                st.markdown("**🔴 Top E-Tractor Signals by Impact**")
                if impact['top_signals']:
                    top_df = pd.DataFrame([
                        {
                            'Title': s['title'],
                            'Dimension': s.get('primary_dimension', ''),
                            'Severity': s.get('disruption_classification', ''),
                            'Impact': round(s.get('impact_score') or 0, 2),
                            'Novelty': round(s.get('novelty_score') or 0, 2),
                            'Velocity': round(s.get('velocity_score') or 0, 2),
                            'Source URL': s.get('url', ''),
                        }
                        for s in impact['top_signals']
                    ])
                    st.dataframe(
                        top_df, use_container_width=True, hide_index=True,
                        column_config={
                            "Title": st.column_config.TextColumn("Signal", width="large"),
                            "Impact": st.column_config.NumberColumn("Impact", format="%.2f"),
                            "Novelty": st.column_config.NumberColumn("Novelty", format="%.2f"),
                            "Velocity": st.column_config.NumberColumn("Velocity", format="%.2f"),
                            "Source URL": st.column_config.LinkColumn("🔗 Source"),
                        }
                    )

                # Competitor signals
                if competitor_signals:
                    st.markdown(f"**🏭 {len(competitor_signals)} Competitor-Related Signal(s) Detected**")
                    for cs in competitor_signals[:5]:
                        with st.expander(f"🏭 {cs['title']} [{cs.get('primary_dimension','')}]"):
                            st.caption(cs.get('content', '')[:400] + '...' if len(cs.get('content','')) > 400 else cs.get('content',''))
                            if cs.get('url'):
                                st.markdown(f"[🔗 Source]({cs['url']})")
                else:
                    st.info("No competitor signals in current database. Expand scouting scope to competitor press releases and patent filings.")

        # ── SECTION 3: Sales Recommendations ─────────────────────────
        with section_tabs[2]:
            st.markdown("### Sales Recommendations & Insights")
            st.caption(
                "Structured commercial recommendations for AGCO/Fendt sales leadership, "
                "derived from subsidy program analysis and live signal intelligence."
            )

            sales_recs = etractor.generate_sales_recommendations(signals)
            for rec in sales_recs:
                priority_color = SEVERITY_COLORS.get(rec['priority'], '#888')
                badge = PRIORITY_BADGE.get(rec['priority'], '⚪')
                with st.expander(f"{badge} [{rec['area']}] {rec['title']}", expanded=True):
                    col_left, col_right = st.columns([3, 1])
                    with col_left:
                        st.markdown(f"**Recommendation:**  \n{rec['recommendation']}")
                        st.markdown(f"**Recommended Action:**  \n_{rec['action']}_")
                    with col_right:
                        st.markdown(
                            f"<div style='padding:10px;background:rgba(255,255,255,0.04);border-radius:6px;text-align:center;'>"
                            f"<p style='color:#888;margin:0;font-size:11px;'>KPI</p>"
                            f"<p style='color:#e0e0e0;margin:4px 0;font-size:12px;'>{rec['kpi']}</p>"
                            f"<p style='color:{priority_color};margin:0;font-size:12px;font-weight:bold;'>{rec['priority']}</p>"
                            f"</div>",
                            unsafe_allow_html=True
                        )

        # ── SECTION 4: Marketing Recommendations ─────────────────────
        with section_tabs[3]:
            st.markdown("### Marketing Recommendations & Insights")
            st.caption(
                "Strategy, campaign concepts, and communication measures "
                "tailored to e-tractor adoption barriers and Fendt's market position."
            )

            marketing_recs = etractor.generate_marketing_recommendations()
            for rec in marketing_recs:
                priority_color = SEVERITY_COLORS.get(rec['priority'], '#888')
                badge = PRIORITY_BADGE.get(rec['priority'], '⚪')
                with st.expander(f"{badge} [{rec['category']}] {rec['title']}", expanded=True):
                    st.markdown(f"**Why this matters:**  \n{rec['rationale']}")
                    st.markdown(
                        f"<div style='padding:12px;background:rgba(0,180,80,0.06);border-radius:6px;"
                        f"border-left:3px solid #00cc55;margin:10px 0;'>"
                        f"<p style='color:#aaa;margin:0;font-size:13px;'>"
                        f"<b>💡 Campaign Concept:</b> {rec['campaign_concept']}</p>"
                        f"</div>",
                        unsafe_allow_html=True
                    )
                    st.markdown("**Recommended Channels:**")
                    for ch in rec['channels']:
                        st.markdown(f"- {ch}")
                    st.markdown(
                        f"<p style='color:{priority_color};font-size:12px;margin:4px 0;'>"
                        f"<b>KPI:</b> {rec['kpi']}</p>",
                        unsafe_allow_html=True
                    )

        # ── SECTION 5: Customer Needs ─────────────────────────────────
        with section_tabs[4]:
            st.markdown("### E-Tractor Specific Customer Needs")
            st.caption(
                "How electric tractor customer needs differ from standard diesel tractor usage — "
                "and what this means for Fendt product, service, and sales strategy."
            )

            needs = etractor.generate_etractor_customer_needs()
            for need in needs:
                risk_col = SEVERITY_COLORS.get(need['risk_if_not_addressed'], '#888')
                badge = PRIORITY_BADGE.get(need['risk_if_not_addressed'], '⚪')
                with st.expander(f"{badge} {need['need_area']}", expanded=False):
                    col_d, col_e = st.columns(2)
                    with col_d:
                        st.markdown(
                            f"<div style='padding:10px;background:rgba(255,100,0,0.06);border-radius:6px;height:100%;'>"
                            f"<p style='color:#888;font-size:11px;margin:0 0 4px 0;'>DIESEL BASELINE</p>"
                            f"<p style='color:#ccc;font-size:13px;margin:0;'>{need['diesel_baseline']}</p>"
                            f"</div>",
                            unsafe_allow_html=True
                        )
                    with col_e:
                        st.markdown(
                            f"<div style='padding:10px;background:rgba(0,180,80,0.07);border-radius:6px;height:100%;'>"
                            f"<p style='color:#888;font-size:11px;margin:0 0 4px 0;'>E-TRACTOR REQUIREMENT</p>"
                            f"<p style='color:#ccc;font-size:13px;margin:0;'>{need['etractor_requirement']}</p>"
                            f"</div>",
                            unsafe_allow_html=True
                        )
                    st.markdown(
                        f"<div style='padding:10px;background:rgba(0,140,255,0.07);border-radius:6px;"
                        f"border-left:3px solid #00ccff;margin-top:8px;'>"
                        f"<p style='color:#888;font-size:11px;margin:0 0 4px 0;'>FENDT IMPLICATION</p>"
                        f"<p style='color:#e0e0e0;font-size:13px;margin:0;'>{need['fendt_implication']}</p>"
                        f"</div>",
                        unsafe_allow_html=True
                    )
                    st.markdown(
                        f"<p style='color:{risk_col};font-size:12px;margin:8px 0 0 0;'>"
                        f"<b>Risk if not addressed: {need['risk_if_not_addressed']}</b></p>",
                        unsafe_allow_html=True
                    )

        # ── SECTION 6: Semantic Search ────────────────────────────────
        with section_tabs[5]:
            st.markdown("### 🔍 Semantic Signal Search")
            st.markdown(
                f"""
                <div style='padding: 12px 16px; background: rgba(255,255,255,0.04); border-radius: 8px; margin-bottom: 14px; border: 1px solid rgba(255,255,255,0.08);'>
                    <p style='color: #aaa; margin: 0; font-size: 13px;'>
                    Search the signal database by meaning, not exact keywords.
                    Uses <b>{'TF-IDF vector similarity (scikit-learn)' if searcher.is_sklearn_available() else 'keyword fallback (install scikit-learn for semantic search)'}</b>.
                    Index covers <b>{searcher.index_size()} signals</b>.
                    </p>
                </div>
                """,
                unsafe_allow_html=True
            )

            query_col, k_col = st.columns([4, 1])
            with query_col:
                search_query = st.text_input(
                    "🔍 Enter a topic or question",
                    placeholder="e.g. EU subsidies for electric agricultural machinery 2026",
                    help="Search by concept — the engine finds signals semantically similar to your query."
                )
            with k_col:
                top_k = st.number_input("Results", min_value=3, max_value=20, value=8, step=1)

            preset_queries = [
                "EU subsidies electric tractors",
                "battery technology agriculture",
                "charging infrastructure rural",
                "competitor electric tractor launch",
                "CAP reform green deal farm",
                "EU battery mandate compliance 2027",
            ]
            st.caption("Quick searches:")
            preset_cols = st.columns(len(preset_queries))
            for i, pq in enumerate(preset_queries):
                with preset_cols[i]:
                    if st.button(pq, key=f"preset_{i}"):
                        search_query = pq

            if search_query:
                with st.spinner("Searching signal database..."):
                    results = searcher.search(search_query, top_k=top_k)

                if results:
                    st.markdown(f"**{len(results)} results for: _{search_query}_**")
                    for sig, score in results:
                        severity = sig.get('disruption_classification', 'LOW')
                        badge = PRIORITY_BADGE.get(severity, '⚪')
                        sev_color = SEVERITY_COLORS.get(severity, '#888')
                        with st.expander(
                            f"{badge} {sig['title']} [{sig.get('primary_dimension', '')}] — Score: {score:.3f}"
                        ):
                            col_s1, col_s2 = st.columns([3, 1])
                            with col_s1:
                                st.caption(
                                    sig.get('content', '')[:400] + '...'
                                    if len(sig.get('content', '')) > 400
                                    else sig.get('content', '')
                                )
                                if sig.get('url'):
                                    st.markdown(f"[🔗 Open Source]({sig['url']})")
                            with col_s2:
                                st.metric("Impact", f"{sig.get('impact_score') or 0:.2f}")
                                st.metric("Novelty", f"{sig.get('novelty_score') or 0:.2f}")
                                st.markdown(
                                    f"<p style='color:{sev_color};font-size:12px;margin:4px 0;'>"
                                    f"<b>{severity}</b></p>",
                                    unsafe_allow_html=True
                                )
                else:
                    st.info("No matching signals found. Try a broader search term or run the intelligence sweep to add more signals.")
            else:
                st.info("Enter a search query above to find semantically similar signals in the database.")

# ===========================
# FOOTER
# ===========================

st.sidebar.markdown("---")
st.sidebar.caption("© 2026 Fendt Strategic Intelligence")
st.sidebar.caption("Powered by Claude Code & SQLite")
